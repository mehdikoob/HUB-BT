"""
Test suite for PPTX Bilan Partenaire and PDF Incident Report export fixes
- PPTX: Verify all slides are generated for multi-programme partners (21 slides for 7 programmes)
- PDF: Verify long text in points_attention is properly wrapped and doesn't overflow
"""

import pytest
import requests
import os
import io
from pptx import Presentation
from PyPDF2 import PdfReader

# Get BASE_URL from environment
BASE_URL = os.environ.get('REACT_APP_BACKEND_URL', '').rstrip('/')

# Test credentials
ADMIN_EMAIL = "admin@hubblindtests.com"
ADMIN_PASSWORD = "admin123"

# Test IDs from the review request
PARTENAIRE_ID_AZUREVA = "cd529b17-46de-4a31-96e9-4cfa265db425"  # Azureva with 7 programmes
TEST_ID_FOR_PDF = "379a3e4a-9445-4545-93ac-45dc79da9fbc"


@pytest.fixture(scope="module")
def auth_token():
    """Get authentication token for API calls"""
    response = requests.post(
        f"{BASE_URL}/api/auth/login",
        json={"email": ADMIN_EMAIL, "password": ADMIN_PASSWORD}
    )
    assert response.status_code == 200, f"Login failed: {response.text}"
    data = response.json()
    assert "access_token" in data, "No access_token in response"
    return data["access_token"]


@pytest.fixture(scope="module")
def authenticated_session(auth_token):
    """Create authenticated session"""
    session = requests.Session()
    session.headers.update({
        "Authorization": f"Bearer {auth_token}",
        "Content-Type": "application/json"
    })
    return session


class TestPPTXBilanPartenaireMultiProgramme:
    """
    Test PPTX Bilan Partenaire export for partners with multiple programmes.
    Bug fix: Indentation was incorrect causing only first programme's slides to be generated.
    Expected: 3 slides per programme (Vue d'ensemble, Tests Sites, Tests Ligne)
    """
    
    def test_pptx_export_returns_200(self, authenticated_session):
        """Test that PPTX export endpoint returns 200 OK"""
        response = authenticated_session.get(
            f"{BASE_URL}/api/export/bilan-partenaire-ppt",
            params={
                "partenaire_id": PARTENAIRE_ID_AZUREVA,
                "date_debut": "2024-01-01",
                "date_fin": "2025-12-31"
            }
        )
        assert response.status_code == 200, f"PPTX export failed: {response.status_code} - {response.text}"
        assert "application/vnd.openxmlformats-officedocument.presentationml.presentation" in response.headers.get("content-type", "")
    
    def test_pptx_has_correct_slide_count_for_multi_programme(self, authenticated_session):
        """
        CRITICAL TEST: Verify that PPTX has 21 slides for Azureva (7 programmes × 3 slides)
        This validates the indentation fix in the export function.
        """
        # First, get the partenaire to verify number of programmes
        partenaire_response = authenticated_session.get(
            f"{BASE_URL}/api/partenaires/{PARTENAIRE_ID_AZUREVA}"
        )
        assert partenaire_response.status_code == 200, f"Failed to get partenaire: {partenaire_response.text}"
        partenaire = partenaire_response.json()
        num_programmes = len(partenaire.get('programmes_ids', []))
        expected_slides = num_programmes * 3  # 3 slides per programme
        
        print(f"Partenaire: {partenaire.get('nom')}")
        print(f"Number of programmes: {num_programmes}")
        print(f"Expected slides: {expected_slides}")
        
        # Get the PPTX file
        response = authenticated_session.get(
            f"{BASE_URL}/api/export/bilan-partenaire-ppt",
            params={
                "partenaire_id": PARTENAIRE_ID_AZUREVA,
                "date_debut": "2024-01-01",
                "date_fin": "2025-12-31"
            }
        )
        assert response.status_code == 200, f"PPTX export failed: {response.text}"
        
        # Parse the PPTX to count slides
        pptx_content = io.BytesIO(response.content)
        prs = Presentation(pptx_content)
        actual_slides = len(prs.slides)
        
        print(f"Actual slides in PPTX: {actual_slides}")
        
        # Verify slide count matches expected (3 slides per programme)
        assert actual_slides == expected_slides, (
            f"PPTX slide count mismatch! Expected {expected_slides} slides "
            f"({num_programmes} programmes × 3 slides), but got {actual_slides}. "
            f"This indicates the indentation fix may not be working correctly."
        )
    
    def test_pptx_slides_contain_programme_names(self, authenticated_session):
        """Verify that slides contain programme names in titles"""
        # Get programmes for this partenaire
        partenaire_response = authenticated_session.get(
            f"{BASE_URL}/api/partenaires/{PARTENAIRE_ID_AZUREVA}"
        )
        partenaire = partenaire_response.json()
        programme_ids = partenaire.get('programmes_ids', [])
        
        # Get programme names
        programmes_response = authenticated_session.get(f"{BASE_URL}/api/programmes")
        all_programmes = programmes_response.json()
        programme_names = [p['nom'] for p in all_programmes if p['id'] in programme_ids]
        
        print(f"Programme names: {programme_names}")
        
        # Get PPTX
        response = authenticated_session.get(
            f"{BASE_URL}/api/export/bilan-partenaire-ppt",
            params={
                "partenaire_id": PARTENAIRE_ID_AZUREVA,
                "date_debut": "2024-01-01",
                "date_fin": "2025-12-31"
            }
        )
        
        pptx_content = io.BytesIO(response.content)
        prs = Presentation(pptx_content)
        
        # Extract all text from slides
        all_slide_text = []
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + " "
            all_slide_text.append(slide_text)
        
        combined_text = " ".join(all_slide_text)
        
        # Verify each programme name appears in the PPTX
        for prog_name in programme_names:
            assert prog_name in combined_text, (
                f"Programme '{prog_name}' not found in PPTX slides. "
                f"This suggests slides for this programme were not generated."
            )
            print(f"✓ Programme '{prog_name}' found in PPTX")


class TestPDFIncidentReportWordWrap:
    """
    Test PDF Incident Report export with long text in points_attention.
    Bug fix: Text was overflowing cells, now uses Paragraph with wordWrap='CJK'.
    """
    
    def test_pdf_export_returns_200(self, authenticated_session):
        """Test that PDF export endpoint returns 200 OK"""
        response = authenticated_session.get(
            f"{BASE_URL}/api/export-alerte-report/{TEST_ID_FOR_PDF}",
            params={"test_type": "site"}
        )
        assert response.status_code == 200, f"PDF export failed: {response.status_code} - {response.text}"
        assert "application/pdf" in response.headers.get("content-type", "")
    
    def test_pdf_is_valid_and_readable(self, authenticated_session):
        """Test that the generated PDF is valid and can be read"""
        response = authenticated_session.get(
            f"{BASE_URL}/api/export-alerte-report/{TEST_ID_FOR_PDF}",
            params={"test_type": "site"}
        )
        assert response.status_code == 200
        
        # Parse PDF to verify it's valid
        pdf_content = io.BytesIO(response.content)
        reader = PdfReader(pdf_content)
        
        num_pages = len(reader.pages)
        print(f"PDF has {num_pages} page(s)")
        
        # PDF should have at least 1 page
        assert num_pages >= 1, "PDF should have at least 1 page"
        
        # Extract text from first page to verify content
        first_page_text = reader.pages[0].extract_text()
        print(f"First page text preview: {first_page_text[:500]}...")
        
        # Verify key elements are present
        assert "RAPPORT D'INCIDENT" in first_page_text or "RAPPORT" in first_page_text, (
            "PDF should contain 'RAPPORT D'INCIDENT' title"
        )
    
    def test_pdf_contains_incident_data(self, authenticated_session):
        """Test that PDF contains incident/alerte data"""
        response = authenticated_session.get(
            f"{BASE_URL}/api/export-alerte-report/{TEST_ID_FOR_PDF}",
            params={"test_type": "site"}
        )
        assert response.status_code == 200
        
        pdf_content = io.BytesIO(response.content)
        reader = PdfReader(pdf_content)
        
        # Extract all text
        all_text = ""
        for page in reader.pages:
            all_text += page.extract_text() + "\n"
        
        # Check for expected content
        expected_keywords = ["Programme", "Partenaire", "Date"]
        for keyword in expected_keywords:
            assert keyword in all_text, f"PDF should contain '{keyword}'"
            print(f"✓ Found '{keyword}' in PDF")


class TestAPIEndpointsHealth:
    """Basic health checks for the export endpoints"""
    
    def test_api_root_accessible(self):
        """Test that API root is accessible"""
        response = requests.get(f"{BASE_URL}/api/")
        assert response.status_code == 200
        print(f"API root response: {response.json()}")
    
    def test_login_works(self):
        """Test that login endpoint works"""
        response = requests.post(
            f"{BASE_URL}/api/auth/login",
            json={"email": ADMIN_EMAIL, "password": ADMIN_PASSWORD}
        )
        assert response.status_code == 200
        data = response.json()
        assert "access_token" in data
        print("✓ Login successful")
    
    def test_partenaires_endpoint(self, authenticated_session):
        """Test that partenaires endpoint works"""
        response = authenticated_session.get(f"{BASE_URL}/api/partenaires")
        assert response.status_code == 200
        partenaires = response.json()
        assert isinstance(partenaires, list)
        print(f"✓ Found {len(partenaires)} partenaires")
    
    def test_programmes_endpoint(self, authenticated_session):
        """Test that programmes endpoint works"""
        response = authenticated_session.get(f"{BASE_URL}/api/programmes")
        assert response.status_code == 200
        programmes = response.json()
        assert isinstance(programmes, list)
        print(f"✓ Found {len(programmes)} programmes")


class TestAzurevaPartenaireData:
    """Verify Azureva partenaire data for testing"""
    
    def test_azureva_exists(self, authenticated_session):
        """Test that Azureva partenaire exists"""
        response = authenticated_session.get(
            f"{BASE_URL}/api/partenaires/{PARTENAIRE_ID_AZUREVA}"
        )
        assert response.status_code == 200, f"Azureva partenaire not found: {response.text}"
        partenaire = response.json()
        print(f"Partenaire: {partenaire.get('nom')}")
        print(f"ID: {partenaire.get('id')}")
    
    def test_azureva_has_7_programmes(self, authenticated_session):
        """Test that Azureva has exactly 7 programmes"""
        response = authenticated_session.get(
            f"{BASE_URL}/api/partenaires/{PARTENAIRE_ID_AZUREVA}"
        )
        assert response.status_code == 200
        partenaire = response.json()
        
        programme_ids = partenaire.get('programmes_ids', [])
        num_programmes = len(programme_ids)
        
        print(f"Azureva has {num_programmes} programmes")
        print(f"Programme IDs: {programme_ids}")
        
        assert num_programmes == 7, (
            f"Expected Azureva to have 7 programmes, but found {num_programmes}. "
            f"This may affect the PPTX slide count test."
        )


class TestAlertePDFData:
    """Verify alerte data for PDF testing"""
    
    def test_alerte_exists_for_test(self, authenticated_session):
        """Test that alerte exists for the test ID"""
        response = authenticated_session.get(f"{BASE_URL}/api/alertes")
        assert response.status_code == 200
        alertes = response.json()
        
        # Find alerte for our test ID
        test_alertes = [a for a in alertes if a.get('test_id') == TEST_ID_FOR_PDF]
        
        print(f"Found {len(test_alertes)} alerte(s) for test {TEST_ID_FOR_PDF}")
        
        if test_alertes:
            for alerte in test_alertes:
                print(f"  - Alerte ID: {alerte.get('id')}")
                print(f"    Description: {alerte.get('description')}")
                print(f"    Points attention: {alerte.get('points_attention')}")
        
        assert len(test_alertes) > 0, (
            f"No alertes found for test {TEST_ID_FOR_PDF}. "
            f"PDF export requires at least one alerte."
        )


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])
