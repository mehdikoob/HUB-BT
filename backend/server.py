from fastapi import FastAPI, APIRouter, HTTPException, Query, UploadFile, File, Depends
from fastapi.responses import StreamingResponse, FileResponse
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.staticfiles import StaticFiles
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
import csv
import io
import shutil
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict, field_validator
from typing import List, Optional
import uuid
from datetime import datetime, timezone, time
from enum import Enum
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy
import zipfile
import shutil
from lxml import etree
from passlib.context import CryptContext
from jose import JWTError, jwt
from datetime import timedelta
import calendar
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image
from reportlab.lib.enums import TA_CENTER, TA_LEFT

ROOT_DIR = Path(__file__).parent
UPLOAD_DIR = ROOT_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)
TEMPLATE_DIR = ROOT_DIR / "templates"
TEMPLATE_DIR.mkdir(exist_ok=True)
load_dotenv(ROOT_DIR / '.env')

# JWT and Password hashing configuration
SECRET_KEY = os.getenv('SECRET_KEY', 'your-secret-key-change-in-production-please')
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 24 * 7  # 7 days

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/api/auth/login", auto_error=False)

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Create the main app without a prefix
app = FastAPI()

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# Enums
class StatutAlerte(str, Enum):
    ouvert = "ouvert"
    resolu = "resolu"

class TypeTest(str, Enum):
    TS = "TS"
    TL = "TL"

class EvaluationAccueil(str, Enum):
    excellent = "Excellent"
    bien = "Bien"
    moyen = "Moyen"
    mediocre = "Médiocre"

# Models - Programme
class ProgrammeBase(BaseModel):
    nom: str
    description: Optional[str] = None
    logo_url: Optional[str] = None
    url_plateforme: Optional[str] = None  # URL de la plateforme du programme
    identifiant: Optional[str] = None  # Identifiant de connexion
    mot_de_passe: Optional[str] = None  # Mot de passe

class ProgrammeCreate(ProgrammeBase):
    pass

class Programme(ProgrammeBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Models - Partenaire
class ProgrammeContact(BaseModel):
    programme_id: str
    url_site: Optional[str] = None
    code_promo: Optional[str] = None
    numero_telephone: Optional[str] = None
    referer: Optional[str] = None  # URL Referer pour les tests
    test_site_requis: bool = True  # Test site requis pour ce programme
    test_ligne_requis: bool = True  # Test ligne requis pour ce programme

class PartenaireBase(BaseModel):
    nom: str
    programmes_ids: List[str] = []
    contacts_programmes: List[ProgrammeContact] = []  # URLs et téléphones par programme
    naming_attendu: Optional[str] = None
    remise_minimum: Optional[float] = None  # Remise minimum attendue en %
    logo_url: Optional[str] = None
    contact_email: Optional[str] = None  # Email du contact principal

class PartenaireCreate(PartenaireBase):
    pass

class Partenaire(PartenaireBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Models - Test Site
class TestSiteBase(BaseModel):
    programme_id: str
    partenaire_id: str
    date_test: datetime
    application_remise: bool
    prix_public: float
    prix_remise: float
    naming_constate: Optional[str] = None
    cumul_codes: bool
    commentaire: Optional[str] = None
    attachments: List[str] = []  # URLs ou chemins des fichiers joints (jpeg, png, pdf)

class TestSiteCreate(TestSiteBase):
    pass

class CreatedByInfo(BaseModel):
    """Informations sur l'utilisateur créateur"""
    id: str
    nom: str
    prenom: str
    email: str
    role: str

class TestSite(TestSiteBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    pct_remise_calcule: float = 0.0
    user_id: Optional[str] = None  # ID de l'utilisateur qui a créé le test
    created_by: Optional[CreatedByInfo] = None  # Informations du créateur
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Models - Test Ligne
class TestLigneBase(BaseModel):
    programme_id: str
    partenaire_id: str
    date_test: datetime
    numero_telephone: str
    messagerie_vocale_dediee: bool
    decroche_dedie: bool
    delai_attente: str  # Format mm:ss
    nom_conseiller: Optional[str] = "NC"
    evaluation_accueil: EvaluationAccueil
    application_offre: bool
    commentaire: Optional[str] = None

    @field_validator('delai_attente')
    def validate_delai(cls, v):
        parts = v.split(':')
        if len(parts) != 2:
            raise ValueError('Format invalide, attendu mm:ss')
        try:
            minutes = int(parts[0])
            seconds = int(parts[1])
            if minutes < 0 or seconds < 0 or seconds > 59:
                raise ValueError('Valeurs invalides')
            if minutes >= 10:
                raise ValueError('Délai d\'attente trop long (>= 10 min)')
        except ValueError as e:
            raise ValueError(f'Format de délai invalide: {e}')
        return v

class TestLigneCreate(TestLigneBase):
    pass

class TestLigne(TestLigneBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    user_id: Optional[str] = None  # ID de l'utilisateur qui a créé le test
    created_by: Optional[CreatedByInfo] = None  # Informations du créateur
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Models - Alerte
class AlerteBase(BaseModel):
    test_id: str
    type_test: TypeTest
    description: str
    statut: StatutAlerte = StatutAlerte.ouvert
    programme_id: Optional[str] = None
    partenaire_id: Optional[str] = None

class AlerteCreate(AlerteBase):
    pass

class Alerte(AlerteBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    user_id: Optional[str] = None  # ID de l'utilisateur qui a créé/détecté l'alerte
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    resolved_at: Optional[datetime] = None

# Models - Notification
class NotificationBase(BaseModel):
    user_id: str  # Chef de projet qui reçoit la notification
    alerte_id: str  # ID de l'alerte concernée
    programme_id: str  # Programme concerné
    partenaire_id: str  # Partenaire concerné
    message: str  # Message de la notification
    read: bool = False  # Notification lue ou non

class NotificationCreate(NotificationBase):
    pass

class Notification(NotificationBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Models - Email Template
class EmailTemplateBase(BaseModel):
    name: str
    subject_template: str
    body_template: str
    is_default: bool = False

class EmailTemplateCreate(EmailTemplateBase):
    pass

class EmailTemplate(EmailTemplateBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Models - User Signature
class UserSignatureBase(BaseModel):
    user_name: str
    signature_text: str

class UserSignatureCreate(UserSignatureBase):
    pass

class UserSignature(UserSignatureBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Models - Email Draft
class EmailDraftStatus(str, Enum):
    draft = "draft"
    sent = "sent"

class EmailDraftBase(BaseModel):
    alerte_id: str
    template_id: Optional[str] = None
    subject: str
    body: str
    recipient: str
    status: EmailDraftStatus = EmailDraftStatus.draft

class EmailDraftCreate(EmailDraftBase):
    pass

class EmailDraft(EmailDraftBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    sent_at: Optional[datetime] = None

# Models - Email History
class EmailHistoryBase(BaseModel):
    alerte_id: str
    draft_id: str
    recipient: str
    subject: str
    body: str
    status: str  # 'success' or 'failed'
    error_message: Optional[str] = None

class EmailHistory(EmailHistoryBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    sent_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Models - User & Authentication
class UserRole(str, Enum):
    admin = "admin"
    agent = "agent"
    programme = "programme"
    partenaire = "partenaire"
    chef_projet = "chef_projet"  # Chef de projet avec droits admin, affilié à des programmes

class UserBase(BaseModel):
    email: str
    nom: str
    prenom: str
    role: UserRole = UserRole.agent
    is_active: bool = True
    programme_id: Optional[str] = None  # Pour les utilisateurs de type "programme"
    partenaire_id: Optional[str] = None  # Pour les utilisateurs de type "partenaire"
    programme_ids: List[str] = []  # Pour les chefs de projet (plusieurs programmes)

class UserCreate(UserBase):
    password: str

class UserUpdate(BaseModel):
    nom: Optional[str] = None
    prenom: Optional[str] = None
    role: Optional[UserRole] = None
    is_active: Optional[bool] = None
    password: Optional[str] = None
    programme_id: Optional[str] = None
    partenaire_id: Optional[str] = None
    programme_ids: Optional[List[str]] = None  # Pour les chefs de projet

class User(UserBase):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class UserInDB(User):
    password_hash: str

class Token(BaseModel):
    access_token: str
    token_type: str

class TokenData(BaseModel):
    email: Optional[str] = None

class LoginRequest(BaseModel):
    email: str
    password: str

# Helper functions
def calculate_remise_percentage(prix_public: float, prix_remise: float) -> float:
    if prix_public <= 0:
        return 0.0
    return round((1 - prix_remise / prix_public) * 100, 2)

async def replace_template_variables(template_text: str, alerte_id: str) -> str:
    """Replace template variables with actual data from alerte"""
    # Get alerte data
    alerte = await db.alertes.find_one({"id": alerte_id})
    if not alerte:
        return template_text
    
    # Get related data
    programme = await db.programmes.find_one({"id": alerte.get('programme_id')}) if alerte.get('programme_id') else None
    partenaire = await db.partenaires.find_one({"id": alerte.get('partenaire_id')}) if alerte.get('partenaire_id') else None
    
    # Get test data
    test = None
    if alerte['type_test'] == 'TS':
        test = await db.tests_site.find_one({"id": alerte['test_id']})
    else:
        test = await db.tests_ligne.find_one({"id": alerte['test_id']})
    
    # Build replacement dictionary
    replacements = {
        '[Nom du programme]': programme['nom'] if programme else 'N/A',
        '[Nature du problème constaté]': alerte.get('description', 'N/A'),
        '[Date du test]': datetime.fromisoformat(test['date_test']).strftime('%d/%m/%Y') if test and test.get('date_test') else 'N/A',
        '[Nom du site / canal du test]': 'Site web' if alerte['type_test'] == 'TS' else 'Téléphone',
        '[Remise attendue]': f"{partenaire.get('remise_minimum', 'N/A')} %" if partenaire else 'N/A',
        '[Observation]': alerte.get('description', 'N/A'),
        '[Nom du contact]': partenaire.get('contact_email', 'N/A') if partenaire else 'N/A',
    }
    
    # Replace variables
    for key, value in replacements.items():
        template_text = template_text.replace(key, str(value))
    
    return template_text

async def send_email_smtp(recipient: str, subject: str, body: str, signature: str = "") -> dict:
    """Send email via SMTP Outlook"""
    try:
        # Get SMTP configuration from environment
        smtp_server = os.getenv('SMTP_SERVER', 'smtp.office365.com')
        smtp_port = int(os.getenv('SMTP_PORT', '587'))
        smtp_user = os.getenv('SMTP_USER', 'automatisation@qwertys.fr')
        smtp_password = os.getenv('SMTP_PASSWORD', '')
        
        if not smtp_password:
            return {"status": "error", "message": "SMTP configuration not available"}
        
        # Create message
        msg = MIMEMultipart()
        msg['From'] = smtp_user
        msg['To'] = recipient
        msg['Subject'] = subject
        
        # Add body with signature
        full_body = body
        if signature:
            full_body += f"\n\n{signature}"
        
        msg.attach(MIMEText(full_body, 'plain'))
        
        # Send email
        server = smtplib.SMTP(smtp_server, smtp_port)
        server.starttls()
        server.login(smtp_user, smtp_password)
        server.send_message(msg)
        server.quit()
        
        return {"status": "success", "message": "Email sent successfully"}
    except Exception as e:
        logging.error(f"SMTP Error: {str(e)}")
        return {"status": "error", "message": str(e)}

async def create_email_draft_for_alerte(alerte_id: str):
    """Automatically create an email draft when an alerte is created"""
    try:
        # Get alerte
        alerte = await db.alertes.find_one({"id": alerte_id})
        if not alerte:
            return
        
        # Get partenaire to get recipient email
        partenaire = await db.partenaires.find_one({"id": alerte.get('partenaire_id')}) if alerte.get('partenaire_id') else None
        if not partenaire or not partenaire.get('contact_email'):
            logging.warning(f"No contact email for alerte {alerte_id}")
            return
        
        # Get default template
        default_template = await db.email_templates.find_one({"is_default": True})
        if not default_template:
            # Create default template if it doesn't exist
            default_template = EmailTemplate(
                name="Template par défaut",
                subject_template="[Nom du programme] – [Nature du problème constaté]",
                body_template="""Bonjour,

J'espère que vous allez bien. 

Dans le cadre de nos tests à l'aveugle réalisés régulièrement sur [Nom du site / canal du test], notre équipe a relevé un point qui pourrait nécessiter une vérification de votre côté.

Détails du test :

Date du test : [Date du test]
Programme concerné : [Nom du programme]
Remise attendue : [Remise attendue]
Observation : [Observation]

Il est possible qu'il s'agisse d'un cas isolé ou lié à nos conditions de test. Nous préférons donc vous partager l'information afin que vous puissiez vérifier de votre côté et confirmer si tout fonctionne normalement.

Merci de votre retour,
Bien cordialement,""",
                is_default=True
            )
            doc = default_template.model_dump()
            doc['created_at'] = doc['created_at'].isoformat()
            await db.email_templates.insert_one(doc)
        
        # Replace variables in template
        # Handle both dict (from DB) and EmailTemplate object (newly created)
        if isinstance(default_template, dict):
            subject_template = default_template['subject_template']
            body_template = default_template['body_template']
            template_id = default_template['id']
        else:
            subject_template = default_template.subject_template
            body_template = default_template.body_template
            template_id = default_template.id
            
        subject = await replace_template_variables(subject_template, alerte_id)
        body = await replace_template_variables(body_template, alerte_id)
        
        # Create draft
        draft = EmailDraft(
            alerte_id=alerte_id,
            template_id=template_id,
            subject=subject,
            body=body,
            recipient=partenaire['contact_email'],
            status=EmailDraftStatus.draft
        )
        doc = draft.model_dump()
        doc['created_at'] = doc['created_at'].isoformat()
        await db.email_drafts.insert_one(doc)
        
        logging.info(f"Email draft created for alerte {alerte_id}")
    except Exception as e:
        logging.error(f"Error creating email draft: {str(e)}")

async def check_and_create_alerte(test_id: str, type_test: TypeTest, description: str, programme_id: str = None, partenaire_id: str = None, user_id: str = None):
    alerte = Alerte(
        test_id=test_id,
        type_test=type_test,
        description=description,
        statut=StatutAlerte.ouvert,
        programme_id=programme_id,
        partenaire_id=partenaire_id,
        user_id=user_id
    )
    doc = alerte.model_dump()
    doc['created_at'] = doc['created_at'].isoformat()
    await db.alertes.insert_one(doc)
    
    # Automatically create email draft for this alerte
    await create_email_draft_for_alerte(alerte.id)
    
    # Create notifications for relevant chefs de projet
    if programme_id:
        await create_notifications_for_chefs_projet(alerte.id, programme_id, partenaire_id, description)

async def create_notifications_for_chefs_projet(alerte_id: str, programme_id: str, partenaire_id: str, description: str):
    """Créer des notifications pour les chefs de projet concernés par cette alerte"""
    try:
        # Trouver tous les chefs de projet qui ont ce programme dans leur liste
        chefs_projet = await db.users.find({
            "role": "chef_projet",
            "is_active": True,
            "programme_ids": programme_id
        }, {"_id": 0}).to_list(100)
        
        if not chefs_projet:
            return
        
        # Récupérer les noms du programme et partenaire pour le message
        programme = await db.programmes.find_one({"id": programme_id}, {"_id": 0, "nom": 1})
        partenaire = await db.partenaires.find_one({"id": partenaire_id}, {"_id": 0, "nom": 1})
        
        programme_nom = programme.get('nom') if programme else 'Programme inconnu'
        partenaire_nom = partenaire.get('nom') if partenaire else 'Partenaire inconnu'
        
        # Créer une notification pour chaque chef de projet concerné
        for chef in chefs_projet:
            notification = Notification(
                user_id=chef['id'],
                alerte_id=alerte_id,
                programme_id=programme_id,
                partenaire_id=partenaire_id,
                message=f"[{programme_nom}] - {partenaire_nom} : {description[:100]}",
                read=False
            )
            
            doc = notification.model_dump()
            doc['created_at'] = doc['created_at'].isoformat()
            await db.notifications.insert_one(doc)
        
        print(f"✅ {len(chefs_projet)} notification(s) créée(s) pour l'alerte {alerte_id}")
        
    except Exception as e:
        print(f"❌ Erreur lors de la création des notifications: {str(e)}")

# Authentication helper functions
def verify_password(plain_password, hashed_password):
    """Verify a password against its hash"""
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password):
    """Hash a password"""
    return pwd_context.hash(password)

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    """Create JWT access token"""
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.now(timezone.utc) + expires_delta
    else:
        expire = datetime.now(timezone.utc) + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

async def get_current_user(token: str = Depends(oauth2_scheme)):
    """Get current user from JWT token"""
    if not token:
        return None
    
    credentials_exception = HTTPException(
        status_code=401,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email: str = payload.get("sub")
        if email is None:
            raise credentials_exception
        token_data = TokenData(email=email)
    except JWTError:
        raise credentials_exception
    
    user = await db.users.find_one({"email": token_data.email})
    if user is None:
        raise credentials_exception
    return User(**user)

async def get_current_active_user(current_user: User = Depends(get_current_user)):
    """Get current active user"""
    if not current_user:
        raise HTTPException(status_code=401, detail="Not authenticated")
    if not current_user.is_active:
        raise HTTPException(status_code=400, detail="Inactive user")
    return current_user

def is_admin_or_chef_projet(user: User) -> bool:
    """Vérifier si l'utilisateur a des droits d'admin (admin ou chef de projet)"""
    return user.role in [UserRole.admin, UserRole.chef_projet]

# Routes - Root
@api_router.get("/")
async def root():
    return {"message": "API QWERTYS Blind Tests"}

# Routes - Programmes
@api_router.post("/programmes", response_model=Programme)
async def create_programme(input: ProgrammeCreate):
    programme = Programme(**input.model_dump())
    doc = programme.model_dump()
    doc['created_at'] = doc['created_at'].isoformat()
    await db.programmes.insert_one(doc)
    return programme

@api_router.get("/programmes", response_model=List[Programme])
async def get_programmes():
    programmes = await db.programmes.find({}, {"_id": 0}).to_list(1000)
    for p in programmes:
        if isinstance(p.get('created_at'), str):
            p['created_at'] = datetime.fromisoformat(p['created_at'])
    return programmes

@api_router.get("/programmes/{programme_id}", response_model=Programme)
async def get_programme(programme_id: str):
    programme = await db.programmes.find_one({"id": programme_id}, {"_id": 0})
    if not programme:
        raise HTTPException(status_code=404, detail="Programme non trouvé")
    if isinstance(programme.get('created_at'), str):
        programme['created_at'] = datetime.fromisoformat(programme['created_at'])
    return programme

@api_router.put("/programmes/{programme_id}", response_model=Programme)
async def update_programme(programme_id: str, input: ProgrammeCreate):
    existing = await db.programmes.find_one({"id": programme_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Programme non trouvé")
    
    update_data = input.model_dump()
    await db.programmes.update_one({"id": programme_id}, {"$set": update_data})
    
    updated = await db.programmes.find_one({"id": programme_id}, {"_id": 0})
    if isinstance(updated.get('created_at'), str):
        updated['created_at'] = datetime.fromisoformat(updated['created_at'])
    return updated

@api_router.delete("/programmes/{programme_id}")
async def delete_programme(programme_id: str):
    result = await db.programmes.delete_one({"id": programme_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Programme non trouvé")
    return {"message": "Programme supprimé"}

# Routes - Partenaires
@api_router.post("/partenaires", response_model=Partenaire)
async def create_partenaire(input: PartenaireCreate):
    partenaire = Partenaire(**input.model_dump())
    doc = partenaire.model_dump()
    doc['created_at'] = doc['created_at'].isoformat()
    await db.partenaires.insert_one(doc)
    return partenaire

@api_router.get("/partenaires", response_model=List[Partenaire])
async def get_partenaires():
    partenaires = await db.partenaires.find({}, {"_id": 0}).to_list(1000)
    for p in partenaires:
        if isinstance(p.get('created_at'), str):
            p['created_at'] = datetime.fromisoformat(p['created_at'])
    return partenaires

@api_router.get("/contacts/all")
async def get_all_contacts():
    """Get all partenaire contacts with their associated programmes"""
    partenaires = await db.partenaires.find({}, {"_id": 0}).to_list(1000)
    programmes = await db.programmes.find({}, {"_id": 0}).to_list(1000)
    
    # Create a map for quick programme lookup
    programmes_map = {p['id']: p['nom'] for p in programmes}
    
    contacts = []
    for partenaire in partenaires:
        # Get all programmes for this partenaire
        programme_names = [programmes_map.get(pid, 'Unknown') for pid in partenaire.get('programmes_ids', [])]
        
        contacts.append({
            'partenaire_id': partenaire['id'],
            'partenaire_nom': partenaire['nom'],
            'contact_email': partenaire.get('contact_email'),
            'programmes': ', '.join(programme_names) if programme_names else 'Aucun'
        })
    
    return contacts

@api_router.get("/partenaires/{partenaire_id}", response_model=Partenaire)
async def get_partenaire(partenaire_id: str):
    partenaire = await db.partenaires.find_one({"id": partenaire_id}, {"_id": 0})
    if not partenaire:
        raise HTTPException(status_code=404, detail="Partenaire non trouvé")
    if isinstance(partenaire.get('created_at'), str):
        partenaire['created_at'] = datetime.fromisoformat(partenaire['created_at'])
    return partenaire

@api_router.put("/partenaires/{partenaire_id}", response_model=Partenaire)
async def update_partenaire(partenaire_id: str, input: PartenaireCreate):
    existing = await db.partenaires.find_one({"id": partenaire_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Partenaire non trouvé")
    
    update_data = input.model_dump()
    await db.partenaires.update_one({"id": partenaire_id}, {"$set": update_data})
    
    updated = await db.partenaires.find_one({"id": partenaire_id}, {"_id": 0})
    if isinstance(updated.get('created_at'), str):
        updated['created_at'] = datetime.fromisoformat(updated['created_at'])
    return updated

@api_router.delete("/partenaires/{partenaire_id}")
async def delete_partenaire(partenaire_id: str):
    result = await db.partenaires.delete_one({"id": partenaire_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Partenaire non trouvé")
    return {"message": "Partenaire supprimé"}

# Route pour vérifier la remise par rapport à la remise minimum attendue
@api_router.get("/partenaires/{partenaire_id}/verify-remise")
async def verify_remise(partenaire_id: str, remise_calculee: float):
    partenaire = await db.partenaires.find_one({"id": partenaire_id}, {"_id": 0})
    if not partenaire:
        raise HTTPException(status_code=404, detail="Partenaire non trouvé")
    
    remise_minimum = partenaire.get('remise_minimum')
    if remise_minimum is None:
        return {"conforme": True, "message": "Aucune remise minimum définie"}
    
    if remise_calculee < remise_minimum:
        return {
            "conforme": False,
            "message": f"Remise insuffisante: {remise_calculee}% < {remise_minimum}% attendu",
            "ecart": remise_minimum - remise_calculee
        }
    
    return {"conforme": True, "message": "Remise conforme"}

# Routes - Tests Site
@api_router.post("/tests-site", response_model=TestSite)
async def create_test_site(input: TestSiteCreate, current_user: User = Depends(get_current_active_user)):
    # Bloquer création pour les rôles programme et partenaire (lecture seule)
    if current_user.role in [UserRole.programme, UserRole.partenaire]:
        raise HTTPException(status_code=403, detail="Vous n'avez pas les permissions pour créer des tests")
    
    # Calculate remise percentage
    pct_remise = calculate_remise_percentage(input.prix_public, input.prix_remise)
    
    test_data = input.model_dump()
    test = TestSite(**test_data, pct_remise_calcule=pct_remise, user_id=current_user.id)
    
    # Récupérer le partenaire pour vérifier la remise minimum
    partenaire = await db.partenaires.find_one({"id": input.partenaire_id}, {"_id": 0})
    
    # Validations et création d'alertes
    if input.prix_remise > input.prix_public:
        await check_and_create_alerte(
            test.id,
            TypeTest.TS,
            f"Prix remisé ({input.prix_remise}€) supérieur au prix public ({input.prix_public}€)",
            input.programme_id,
            input.partenaire_id,
            current_user.id
        )
    
    if not input.application_remise:
        await check_and_create_alerte(
            test.id,
            TypeTest.TS,
            "Remise non appliquée",
            input.programme_id,
            input.partenaire_id,
            current_user.id
        )
    
    # Vérifier la remise minimum si définie
    if partenaire and partenaire.get('remise_minimum'):
        remise_minimum = partenaire['remise_minimum']
        if pct_remise < remise_minimum:
            await check_and_create_alerte(
                test.id,
                TypeTest.TS,
                f"Remise insuffisante: {pct_remise}% appliquée, {remise_minimum}% attendue (écart: {remise_minimum - pct_remise}%)",
                input.programme_id,
                input.partenaire_id,
                current_user.id
            )
    
    # Save test
    doc = test.model_dump()
    doc['date_test'] = doc['date_test'].isoformat()
    doc['created_at'] = doc['created_at'].isoformat()
    await db.tests_site.insert_one(doc)
    
    return test

@api_router.get("/tests-site", response_model=List[TestSite])
async def get_tests_site(
    programme_id: Optional[str] = Query(None),
    partenaire_id: Optional[str] = Query(None),
    date_debut: Optional[str] = Query(None),
    date_fin: Optional[str] = Query(None),
    current_user: User = Depends(get_current_active_user)
):
    query = {}
    
    # Filtrage automatique selon le rôle
    if current_user.role == UserRole.programme:
        if not current_user.programme_id:
            raise HTTPException(status_code=403, detail="Aucun programme associé à cet utilisateur")
        query['programme_id'] = current_user.programme_id
    elif current_user.role == UserRole.partenaire:
        if not current_user.partenaire_id:
            raise HTTPException(status_code=403, detail="Aucun partenaire associé à cet utilisateur")
        query['partenaire_id'] = current_user.partenaire_id
    else:
        # Admin et Agent peuvent filtrer manuellement
        if programme_id:
            query['programme_id'] = programme_id
        if partenaire_id:
            query['partenaire_id'] = partenaire_id
    
    if date_debut or date_fin:
        query['date_test'] = {}
        if date_debut:
            query['date_test']['$gte'] = date_debut
        if date_fin:
            query['date_test']['$lte'] = date_fin
    
    tests = await db.tests_site.find(query, {"_id": 0}).to_list(1000)
    
    # Enrichir avec les informations de l'utilisateur créateur
    for t in tests:
        if isinstance(t.get('date_test'), str):
            t['date_test'] = datetime.fromisoformat(t['date_test'])
        if isinstance(t.get('created_at'), str):
            t['created_at'] = datetime.fromisoformat(t['created_at'])
        
        # Ajouter les infos de l'utilisateur qui a créé le test
        if t.get('user_id'):
            user = await db.users.find_one({"id": t['user_id']}, {"_id": 0, "password_hash": 0})
            if user:
                t['created_by'] = {
                    "id": user.get('id'),
                    "nom": user.get('nom'),
                    "prenom": user.get('prenom'),
                    "email": user.get('email'),
                    "role": user.get('role')
                }
    
    return tests

@api_router.get("/tests-site/{test_id}", response_model=TestSite)
async def get_test_site(test_id: str):
    test = await db.tests_site.find_one({"id": test_id}, {"_id": 0})
    if not test:
        raise HTTPException(status_code=404, detail="Test non trouvé")
    if isinstance(test.get('date_test'), str):
        test['date_test'] = datetime.fromisoformat(test['date_test'])
    if isinstance(test.get('created_at'), str):
        test['created_at'] = datetime.fromisoformat(test['created_at'])
    return test

@api_router.put("/tests-site/{test_id}", response_model=TestSite)
async def update_test_site(test_id: str, input: TestSiteCreate):
    # Check if test exists
    existing = await db.tests_site.find_one({"id": test_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Test non trouvé")
    
    # Calculate percentage
    pct_remise = calculate_remise_percentage(input.prix_public, input.prix_remise)
    
    # Prepare update data
    update_data = input.model_dump()
    update_data['pct_remise_calcule'] = pct_remise
    update_data['date_test'] = update_data['date_test'].isoformat()
    
    # Update in database
    await db.tests_site.update_one(
        {"id": test_id},
        {"$set": update_data}
    )
    
    # Get updated document
    updated = await db.tests_site.find_one({"id": test_id}, {"_id": 0})
    if isinstance(updated.get('date_test'), str):
        updated['date_test'] = datetime.fromisoformat(updated['date_test'])
    if isinstance(updated.get('created_at'), str):
        updated['created_at'] = datetime.fromisoformat(updated['created_at'])
    
    return updated

@api_router.delete("/tests-site/{test_id}")
async def delete_test_site(test_id: str):
    result = await db.tests_site.delete_one({"id": test_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Test non trouvé")
    return {"message": "Test supprimé"}

@api_router.post("/upload-attachment")
async def upload_attachment(file: UploadFile = File(...)):
    """Upload a file (jpeg, png, pdf) and return the URL"""
    # Validate file type
    allowed_extensions = [".jpg", ".jpeg", ".png", ".pdf"]
    file_extension = Path(file.filename).suffix.lower()
    
    if file_extension not in allowed_extensions:
        raise HTTPException(
            status_code=400, 
            detail=f"Type de fichier non autorisé. Formats acceptés : {', '.join(allowed_extensions)}"
        )
    
    # Validate file size (max 10MB)
    file.file.seek(0, 2)  # Seek to end
    file_size = file.file.tell()
    file.file.seek(0)  # Reset to start
    
    if file_size > 10 * 1024 * 1024:  # 10MB
        raise HTTPException(status_code=400, detail="Fichier trop volumineux (max 10MB)")
    
    # Generate unique filename
    unique_filename = f"{uuid.uuid4()}{file_extension}"
    file_path = UPLOAD_DIR / unique_filename
    
    # Save file
    try:
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur lors de l'enregistrement du fichier: {str(e)}")
    finally:
        file.file.close()
    
    # Return the URL
    return {
        "filename": file.filename,
        "url": f"/uploads/{unique_filename}",
        "size": file_size
    }

# Routes - Tests Ligne
@api_router.post("/tests-ligne", response_model=TestLigne)
async def create_test_ligne(input: TestLigneCreate, current_user: User = Depends(get_current_active_user)):
    # Bloquer création pour les rôles programme et partenaire (lecture seule)
    if current_user.role in [UserRole.programme, UserRole.partenaire]:
        raise HTTPException(status_code=403, detail="Vous n'avez pas les permissions pour créer des tests")
    
    test = TestLigne(**input.model_dump(), user_id=current_user.id)
    
    # Validations et création d'alertes
    if not input.application_offre:
        await check_and_create_alerte(
            test.id,
            TypeTest.TL,
            "Offre non appliquée",
            input.programme_id,
            input.partenaire_id,
            current_user.id
        )
    
    if not input.messagerie_vocale_dediee and not input.decroche_dedie:
        await check_and_create_alerte(
            test.id,
            TypeTest.TL,
            "Ni messagerie dédiée ni décroche dédié détecté",
            input.programme_id,
            input.partenaire_id,
            current_user.id
        )
    
    # Save test
    doc = test.model_dump()
    doc['date_test'] = doc['date_test'].isoformat()
    doc['created_at'] = doc['created_at'].isoformat()
    await db.tests_ligne.insert_one(doc)
    
    return test

@api_router.get("/tests-ligne", response_model=List[TestLigne])
async def get_tests_ligne(
    programme_id: Optional[str] = Query(None),
    partenaire_id: Optional[str] = Query(None),
    date_debut: Optional[str] = Query(None),
    date_fin: Optional[str] = Query(None),
    current_user: User = Depends(get_current_active_user)
):
    query = {}
    
    # Filtrage automatique selon le rôle
    if current_user.role == UserRole.programme:
        if not current_user.programme_id:
            raise HTTPException(status_code=403, detail="Aucun programme associé à cet utilisateur")
        query['programme_id'] = current_user.programme_id
    elif current_user.role == UserRole.partenaire:
        if not current_user.partenaire_id:
            raise HTTPException(status_code=403, detail="Aucun partenaire associé à cet utilisateur")
        query['partenaire_id'] = current_user.partenaire_id
    else:
        # Admin et Agent peuvent filtrer manuellement
        if programme_id:
            query['programme_id'] = programme_id
        if partenaire_id:
            query['partenaire_id'] = partenaire_id
    
    if date_debut or date_fin:
        query['date_test'] = {}
        if date_debut:
            query['date_test']['$gte'] = date_debut
        if date_fin:
            query['date_test']['$lte'] = date_fin
    
    tests = await db.tests_ligne.find(query, {"_id": 0}).to_list(1000)
    
    # Enrichir avec les informations de l'utilisateur créateur
    for t in tests:
        if isinstance(t.get('date_test'), str):
            t['date_test'] = datetime.fromisoformat(t['date_test'])
        if isinstance(t.get('created_at'), str):
            t['created_at'] = datetime.fromisoformat(t['created_at'])
        
        # Ajouter les infos de l'utilisateur qui a créé le test
        if t.get('user_id'):
            user = await db.users.find_one({"id": t['user_id']}, {"_id": 0, "password_hash": 0})
            if user:
                t['created_by'] = {
                    "id": user.get('id'),
                    "nom": user.get('nom'),
                    "prenom": user.get('prenom'),
                    "email": user.get('email'),
                    "role": user.get('role')
                }
    
    return tests

@api_router.get("/tests-ligne/{test_id}", response_model=TestLigne)
async def get_test_ligne(test_id: str):
    test = await db.tests_ligne.find_one({"id": test_id}, {"_id": 0})
    if not test:
        raise HTTPException(status_code=404, detail="Test non trouvé")
    if isinstance(test.get('date_test'), str):
        test['date_test'] = datetime.fromisoformat(test['date_test'])
    if isinstance(test.get('created_at'), str):
        test['created_at'] = datetime.fromisoformat(test['created_at'])
    return test

@api_router.put("/tests-ligne/{test_id}", response_model=TestLigne)
async def update_test_ligne(test_id: str, input: TestLigneCreate):
    # Check if test exists
    existing = await db.tests_ligne.find_one({"id": test_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Test non trouvé")
    
    # Prepare update data
    update_data = input.model_dump()
    update_data['date_test'] = update_data['date_test'].isoformat()
    if update_data.get('delai_attente') and isinstance(update_data['delai_attente'], time):
        update_data['delai_attente'] = update_data['delai_attente'].strftime('%H:%M:%S')
    
    # Update in database
    await db.tests_ligne.update_one(
        {"id": test_id},
        {"$set": update_data}
    )
    
    # Get updated document
    updated = await db.tests_ligne.find_one({"id": test_id}, {"_id": 0})
    if isinstance(updated.get('date_test'), str):
        updated['date_test'] = datetime.fromisoformat(updated['date_test'])
    if isinstance(updated.get('created_at'), str):
        updated['created_at'] = datetime.fromisoformat(updated['created_at'])
    if isinstance(updated.get('delai_attente'), str):
        updated['delai_attente'] = datetime.strptime(updated['delai_attente'], '%H:%M:%S').time()
    
    return updated

@api_router.delete("/tests-ligne/{test_id}")
async def delete_test_ligne(test_id: str):
    result = await db.tests_ligne.delete_one({"id": test_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Test non trouvé")
    return {"message": "Test supprimé"}

# Routes - Incidents
@api_router.get("/alertes", response_model=List[Alerte])
async def get_incidents(
    statut: Optional[StatutAlerte] = Query(None),
    current_user: User = Depends(get_current_active_user)
):
    query = {}
    
    # Filtrage automatique selon le rôle
    if current_user.role == UserRole.programme:
        if not current_user.programme_id:
            raise HTTPException(status_code=403, detail="Aucun programme associé à cet utilisateur")
        query['programme_id'] = current_user.programme_id
    elif current_user.role == UserRole.partenaire:
        if not current_user.partenaire_id:
            raise HTTPException(status_code=403, detail="Aucun partenaire associé à cet utilisateur")
        query['partenaire_id'] = current_user.partenaire_id
    
    if statut:
        query['statut'] = statut
    
    alertes = await db.alertes.find(query, {"_id": 0}).to_list(1000)
    for i in alertes:
        if isinstance(i.get('created_at'), str):
            i['created_at'] = datetime.fromisoformat(i['created_at'])
        if i.get('resolved_at') and isinstance(i['resolved_at'], str):
            i['resolved_at'] = datetime.fromisoformat(i['resolved_at'])
    return alertes

@api_router.get("/alertes/enriched")
async def get_incidents_enriched(statut: Optional[StatutAlerte] = Query(None)):
    """Get alertes with programme and partenaire details"""
    query = {}
    if statut:
        query['statut'] = statut
    
    alertes = await db.alertes.find(query, {"_id": 0}).to_list(1000)
    
    # Enrich with programme and partenaire data
    for alerte in alertes:
        if isinstance(alerte.get('created_at'), str):
            alerte['created_at'] = datetime.fromisoformat(alerte['created_at'])
        if alerte.get('resolved_at') and isinstance(alerte['resolved_at'], str):
            alerte['resolved_at'] = datetime.fromisoformat(alerte['resolved_at'])
        
        # Get programme details
        if alerte.get('programme_id'):
            programme = await db.programmes.find_one({"id": alerte['programme_id']}, {"_id": 0})
            alerte['programme_nom'] = programme['nom'] if programme else None
        
        # Get partenaire details
        if alerte.get('partenaire_id'):
            partenaire = await db.partenaires.find_one({"id": alerte['partenaire_id']}, {"_id": 0})
            if partenaire:
                alerte['partenaire_nom'] = partenaire['nom']
                alerte['partenaire_contact_email'] = partenaire.get('contact_email')
    
    return alertes

@api_router.put("/alertes/{alerte_id}", response_model=Alerte)
async def resolve_incident(alerte_id: str):
    existing = await db.alertes.find_one({"id": alerte_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Incident non trouvé")
    
    resolved_at = datetime.now(timezone.utc).isoformat()
    await db.alertes.update_one(
        {"id": alerte_id},
        {"$set": {"statut": StatutAlerte.resolu, "resolved_at": resolved_at}}
    )
    
    updated = await db.alertes.find_one({"id": alerte_id}, {"_id": 0})
    if isinstance(updated.get('created_at'), str):
        updated['created_at'] = datetime.fromisoformat(updated['created_at'])
    if updated.get('resolved_at') and isinstance(updated['resolved_at'], str):
        updated['resolved_at'] = datetime.fromisoformat(updated['resolved_at'])
    return updated

@api_router.delete("/alertes/{alerte_id}")
async def delete_incident(alerte_id: str):
    """Delete an alerte (only if resolved)"""
    existing = await db.alertes.find_one({"id": alerte_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Incident non trouvé")
    
    # Verify that the alerte is resolved before allowing deletion
    if existing.get('statut') != StatutAlerte.resolu:
        raise HTTPException(status_code=400, detail="Seuls les alertes résolus peuvent être supprimés")
    
    result = await db.alertes.delete_one({"id": alerte_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Incident non trouvé")
    
    return {"message": "Incident supprimé avec succès"}

# Routes - Notifications
@api_router.get("/notifications", response_model=List[Notification])
async def get_notifications(current_user: User = Depends(get_current_active_user)):
    """Récupérer les notifications de l'utilisateur connecté"""
    notifications = await db.notifications.find(
        {"user_id": current_user.id},
        {"_id": 0}
    ).sort("created_at", -1).to_list(100)
    
    for n in notifications:
        if isinstance(n.get('created_at'), str):
            n['created_at'] = datetime.fromisoformat(n['created_at'])
    
    return notifications

@api_router.get("/notifications/unread-count")
async def get_unread_count(current_user: User = Depends(get_current_active_user)):
    """Compter les notifications non lues"""
    count = await db.notifications.count_documents({
        "user_id": current_user.id,
        "read": False
    })
    return {"count": count}

@api_router.put("/notifications/{notification_id}/read")
async def mark_notification_read(
    notification_id: str,
    current_user: User = Depends(get_current_active_user)
):
    """Marquer une notification comme lue"""
    result = await db.notifications.update_one(
        {"id": notification_id, "user_id": current_user.id},
        {"$set": {"read": True}}
    )
    
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="Notification non trouvée")
    
    return {"message": "Notification marquée comme lue"}

@api_router.put("/notifications/mark-all-read")
async def mark_all_notifications_read(current_user: User = Depends(get_current_active_user)):
    """Marquer toutes les notifications comme lues"""
    result = await db.notifications.update_many(
        {"user_id": current_user.id, "read": False},
        {"$set": {"read": True}}
    )
    
    return {"message": f"{result.modified_count} notification(s) marquée(s) comme lue(s)"}

# Routes - Export Incident Report
@api_router.get("/export-alerte-report/{test_id}")
async def export_incident_report(
    test_id: str,
    test_type: str = Query(..., description="Type de test: 'site' ou 'ligne'"),
    current_user: User = Depends(get_current_active_user)
):
    """Generate a PDF alerte report for a test"""
    
    # Récupérer le test selon son type
    if test_type == "site":
        test = await db.tests_site.find_one({"id": test_id})
        test_label = "Test Site"
    elif test_type == "ligne":
        test = await db.tests_ligne.find_one({"id": test_id})
        test_label = "Test Ligne"
    else:
        raise HTTPException(status_code=400, detail="Type de test invalide")
    
    if not test:
        raise HTTPException(status_code=404, detail="Test non trouvé")
    
    # Récupérer les alertes liés à ce test
    alertes = await db.alertes.find({"test_id": test_id}).to_list(length=None)
    
    if not alertes:
        raise HTTPException(status_code=404, detail="Aucun alerte trouvé pour ce test")
    
    # Récupérer les informations du programme et partenaire
    programme = await db.programmes.find_one({"id": test.get("programme_id")})
    partenaire = await db.partenaires.find_one({"id": test.get("partenaire_id")})
    
    # Créer le PDF en mémoire
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.5*inch, bottomMargin=0.5*inch)
    
    # Styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=18,
        textColor=colors.HexColor('#1e3a8a'),
        spaceAfter=12,
        alignment=TA_CENTER
    )
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=14,
        textColor=colors.HexColor('#1e3a8a'),
        spaceAfter=6,
        spaceBefore=12
    )
    normal_style = styles['Normal']
    
    # Contenu du PDF
    story = []
    
    # Logo
    logo_path = ROOT_DIR / "logo-qwertys.png"
    if logo_path.exists():
        logo = Image(str(logo_path), width=1.5*inch, height=0.6*inch)
        logo.hAlign = 'CENTER'
        story.append(logo)
        story.append(Spacer(1, 0.3*inch))
    
    # Titre
    story.append(Paragraph(f"RAPPORT D'INCIDENT - {test_label.upper()}", title_style))
    story.append(Spacer(1, 0.2*inch))
    
    # Date de génération
    date_generation = datetime.now(timezone.utc).strftime("%d/%m/%Y à %H:%M")
    story.append(Paragraph(f"<b>Date de génération :</b> {date_generation}", normal_style))
    story.append(Spacer(1, 0.3*inch))
    
    # Informations du test
    story.append(Paragraph("INFORMATIONS DU TEST", heading_style))
    
    test_data = [
        ["Programme", programme.get('nom', 'N/A') if programme else 'N/A'],
        ["Partenaire", partenaire.get('nom', 'N/A') if partenaire else 'N/A'],
        ["Date du test", test.get('date_test', 'N/A')],
    ]
    
    if test_type == "site":
        test_data.extend([
            ["Site", test.get('site', 'N/A')],
            ["Prix public", f"{test.get('prix_public', 0)}€"],
            ["Prix remisé", f"{test.get('prix_remise', 0)}€"],
            ["% Remise calculé", f"{test.get('pct_remise_calcule', 0)}%"],
            ["Remise appliquée", "Oui" if test.get('application_remise') else "Non"],
        ])
    else:  # ligne
        test_data.extend([
            ["Offre appliquée", "Oui" if test.get('application_offre') else "Non"],
            ["Messagerie dédiée", "Oui" if test.get('messagerie_vocale_dediee') else "Non"],
            ["Décroché dédié", "Oui" if test.get('decroche_dedie') else "Non"],
        ])
    
    test_table = Table(test_data, colWidths=[2.5*inch, 4*inch])
    test_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#e5e7eb')),
        ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTNAME', (1, 0), (1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
    ]))
    
    story.append(test_table)
    story.append(Spacer(1, 0.3*inch))
    
    # Liste des alertes
    story.append(Paragraph(f"INCIDENTS DÉTECTÉS ({len(alertes)})", heading_style))
    story.append(Spacer(1, 0.1*inch))
    
    for idx, alerte in enumerate(alertes, 1):
        incident_data = [
            [f"Incident #{idx}"],
            ["Description", alerte.get('description', 'N/A')],
            ["Statut", alerte.get('statut', 'N/A').upper()],
            ["Date de création", alerte.get('created_at', 'N/A')[:10] if alerte.get('created_at') else 'N/A'],
        ]
        
        incident_table = Table(incident_data, colWidths=[2.5*inch, 4*inch])
        incident_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#fee2e2')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#991b1b')),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
            ('SPAN', (0, 0), (-1, 0)),
            ('BACKGROUND', (0, 1), (0, -1), colors.HexColor('#f3f4f6')),
            ('FONTNAME', (0, 1), (0, -1), 'Helvetica-Bold'),
            ('FONTNAME', (1, 1), (1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 6),
            ('RIGHTPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ]))
        
        story.append(incident_table)
        story.append(Spacer(1, 0.15*inch))
    
    # Footer
    story.append(Spacer(1, 0.3*inch))
    footer_text = f"""
    <para alignment="center">
    <font size="8" color="#6b7280">
    Ce rapport a été généré automatiquement par le système HUB BLIND TESTS QWERTYS<br/>
    Pour toute question, veuillez contacter l'équipe responsable.
    </font>
    </para>
    """
    story.append(Paragraph(footer_text, normal_style))
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    
    # Return as downloadable file
    filename = f"rapport_incident_{test_type}_{test_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
    
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

# Routes - Stats & Dashboard

async def get_agent_dashboard_stats(user: User):
    """Dashboard simplifié pour les agents - focus sur les tâches à faire"""
    from datetime import datetime, timezone
    
    now = datetime.now(timezone.utc)
    year = now.year
    month = now.month
    first_day = datetime(year, month, 1, tzinfo=timezone.utc).isoformat()
    last_day_num = calendar.monthrange(year, month)[1]
    last_day = datetime(year, month, last_day_num, 23, 59, 59, tzinfo=timezone.utc).isoformat()
    
    # Récupérer tous les partenaires et programmes
    partenaires = await db.partenaires.find({}, {"_id": 0}).to_list(1000)
    programmes = await db.programmes.find({}, {"_id": 0}).to_list(1000)
    programmes_dict = {p['id']: p['nom'] for p in programmes}
    
    # Tâches à effectuer : tests manquants ce mois
    taches_tests = []
    
    for partenaire in partenaires:
        part_id = partenaire['id']
        part_nom = partenaire['nom']
        contacts_programmes = partenaire.get('contacts_programmes', [])
        
        for contact in contacts_programmes:
            prog_id = contact.get('programme_id')
            prog_nom = programmes_dict.get(prog_id, 'Programme inconnu')
            test_site_requis = contact.get('test_site_requis', True)
            test_ligne_requis = contact.get('test_ligne_requis', True)
            
            # Vérifier test site ce mois (uniquement si requis)
            if test_site_requis:
                test_site_count = await db.tests_site.count_documents({
                    "partenaire_id": part_id,
                    "programme_id": prog_id,
                    "date_test": {"$gte": first_day, "$lte": last_day}
                })
                
                # Collecter les tests site manquants comme "tâches à faire"
                if test_site_count == 0:
                    taches_tests.append({
                        "partenaire_id": part_id,
                        "partenaire_nom": part_nom,
                        "programme_id": prog_id,
                        "programme_nom": prog_nom,
                        "type_test": "Site",
                        "priorite": "normale"
                    })
            
            # Vérifier test ligne ce mois (uniquement si requis)
            if test_ligne_requis:
                test_ligne_count = await db.tests_ligne.count_documents({
                    "partenaire_id": part_id,
                    "programme_id": prog_id,
                    "date_test": {"$gte": first_day, "$lte": last_day}
                })
                
                # Collecter les tests ligne manquants comme "tâches à faire"
                if test_ligne_count == 0:
                    taches_tests.append({
                        "partenaire_id": part_id,
                        "partenaire_nom": part_nom,
                        "programme_id": prog_id,
                        "programme_nom": prog_nom,
                        "type_test": "Ligne",
                        "priorite": "normale"
                    })
    
    # Incidents en cours (ouverts uniquement)
    incidents_en_cours = await db.alertes.find(
        {"statut": "ouvert"},
        {"_id": 0}
    ).to_list(1000)
    
    # Enrichir les alertes avec les noms de partenaires et programmes
    for alerte in incidents_en_cours:
        # Récupérer le partenaire
        partenaire = await db.partenaires.find_one(
            {"id": alerte.get("partenaire_id")},
            {"_id": 0, "nom": 1}
        )
        alerte["partenaire_nom"] = partenaire.get("nom") if partenaire else "Inconnu"
        
        # Récupérer le programme
        programme = await db.programmes.find_one(
            {"id": alerte.get("programme_id")},
            {"_id": 0, "nom": 1}
        )
        alerte["programme_nom"] = programme.get("nom") if programme else "Inconnu"
    
    # Compter les tests effectués ce mois (pour message encourageant)
    tests_effectues_mois = await db.tests_site.count_documents({
        "date_test": {"$gte": first_day, "$lte": last_day}
    }) + await db.tests_ligne.count_documents({
        "date_test": {"$gte": first_day, "$lte": last_day}
    })
    
    return {
        "role": "agent",
        "taches_tests": taches_tests,
        "total_taches": len(taches_tests),
        "incidents_en_cours": incidents_en_cours,
        "total_incidents": len(incidents_en_cours),
        "tests_effectues_mois": tests_effectues_mois,
        "current_month": month,
        "current_year": year,
        "message_encourageant": get_encouragement_message(tests_effectues_mois)
    }

def get_encouragement_message(tests_count):
    """Génère un message encourageant basé sur le nombre de tests effectués"""
    if tests_count == 0:
        return "C'est parti pour un nouveau mois ! 💪"
    elif tests_count < 10:
        return f"Bon début ! {tests_count} test{'s' if tests_count > 1 else ''} effectué{'s' if tests_count > 1 else ''} ce mois-ci 🎯"
    elif tests_count < 30:
        return f"Beau travail ! {tests_count} tests effectués ce mois-ci 🌟"
    elif tests_count < 60:
        return f"Excellent rythme ! {tests_count} tests effectués ce mois-ci 🚀"
    else:
        return f"Performance exceptionnelle ! {tests_count} tests effectués ce mois-ci 🏆"


@api_router.get("/stats/dashboard")
async def get_dashboard_stats(current_user: User = Depends(get_current_user)):
    from datetime import datetime, timezone
    
    # Si l'utilisateur est un agent, retourner un dashboard simplifié
    if current_user and current_user.role == "agent":
        return await get_agent_dashboard_stats(current_user)
    
    # Pour les autres rôles (admin, programme, partenaire), le dashboard normal
    total_programmes = await db.programmes.count_documents({})
    total_partenaires = await db.partenaires.count_documents({})
    total_incidents_ouverts = await db.alertes.count_documents({"statut": "ouvert"})
    
    # Calculer les dates du mois en cours
    now = datetime.now(timezone.utc)
    year = now.year
    month = now.month
    first_day = datetime(year, month, 1, tzinfo=timezone.utc).isoformat()
    last_day_num = calendar.monthrange(year, month)[1]
    last_day = datetime(year, month, last_day_num, 23, 59, 59, tzinfo=timezone.utc).isoformat()
    
    # Calculer J-5
    days_until_end = last_day_num - now.day
    is_j5_alert = days_until_end <= 5
    
    # Récupérer tous les partenaires et programmes
    partenaires = await db.partenaires.find({}, {"_id": 0}).to_list(1000)
    programmes = await db.programmes.find({}, {"_id": 0}).to_list(1000)
    programmes_dict = {p['id']: p['nom'] for p in programmes}
    
    # Calculer les tests attendus et effectués
    # Pour chaque partenaire x programme : 1 test site + 1 test ligne attendu
    tests_attendus = 0
    tests_effectues = 0
    tests_manquants = []
    
    for partenaire in partenaires:
        part_id = partenaire['id']
        part_nom = partenaire['nom']
        contacts_programmes = partenaire.get('contacts_programmes', [])
        
        # Pour chaque programme associé à ce partenaire
        for contact in contacts_programmes:
            prog_id = contact.get('programme_id')
            prog_nom = programmes_dict.get(prog_id, 'Programme inconnu')
            test_site_requis = contact.get('test_site_requis', True)
            test_ligne_requis = contact.get('test_ligne_requis', True)
            
            # Calculer les tests attendus selon la configuration du programme
            tests_attendus_programme = 0
            if test_site_requis:
                tests_attendus_programme += 1
            if test_ligne_requis:
                tests_attendus_programme += 1
            tests_attendus += tests_attendus_programme
            
            # Vérifier test site ce mois (uniquement si requis)
            test_site_count = 0
            if test_site_requis:
                test_site_count = await db.tests_site.count_documents({
                    "partenaire_id": part_id,
                    "programme_id": prog_id,
                    "date_test": {"$gte": first_day, "$lte": last_day}
                })
                if test_site_count > 0:
                    tests_effectues += 1
            
            # Vérifier test ligne ce mois (uniquement si requis)
            test_ligne_count = 0
            if test_ligne_requis:
                test_ligne_count = await db.tests_ligne.count_documents({
                    "partenaire_id": part_id,
                    "programme_id": prog_id,
                    "date_test": {"$gte": first_day, "$lte": last_day}
                })
                if test_ligne_count > 0:
                    tests_effectues += 1
            
            # Collecter les tests manquants
            manquants = []
            if test_site_requis and test_site_count == 0:
                manquants.append("Site")
            if test_ligne_requis and test_ligne_count == 0:
                manquants.append("Ligne")
            
            if manquants:
                tests_manquants.append({
                    "partenaire_id": part_id,
                    "partenaire_nom": part_nom,
                    "programme_id": prog_id,
                    "programme_nom": prog_nom,
                    "types_manquants": manquants
                })
    
    # Compter le nombre de partenaires uniques avec tests manquants
    partenaires_manquants = len(set([t['partenaire_id'] for t in tests_manquants]))
    
    # Si on est à J-5, compter les partenaires avec tests manquants comme critiques
    tests_manquants_j5 = 0
    if is_j5_alert:
        tests_manquants_j5 = partenaires_manquants
    
    # Taux de réussite TS (sur le mois)
    total_tests_site_mois = await db.tests_site.count_documents({
        "date_test": {"$gte": first_day, "$lte": last_day}
    })
    tests_site_reussis = await db.tests_site.count_documents({
        "date_test": {"$gte": first_day, "$lte": last_day},
        "application_remise": True
    })
    taux_reussite_ts = (tests_site_reussis / total_tests_site_mois * 100) if total_tests_site_mois > 0 else 0
    
    # Taux de réussite TL (sur le mois)
    total_tests_ligne_mois = await db.tests_ligne.count_documents({
        "date_test": {"$gte": first_day, "$lte": last_day}
    })
    tests_ligne_reussis = await db.tests_ligne.count_documents({
        "date_test": {"$gte": first_day, "$lte": last_day},
        "application_offre": True
    })
    taux_reussite_tl = (tests_ligne_reussis / total_tests_ligne_mois * 100) if total_tests_ligne_mois > 0 else 0
    
    # Calcul du nombre réel de tests manquants
    tests_manquants_reel = tests_attendus - tests_effectues
    
    # Calcul de l'indicateur de retard basé sur le prorata du mois
    jour_actuel = now.day
    pourcentage_mois_ecoule = (jour_actuel / last_day_num) * 100
    pourcentage_tests_effectues = (tests_effectues / tests_attendus * 100) if tests_attendus > 0 else 0
    
    # Calculer le retard : différence entre ce qui devrait être fait et ce qui est fait
    retard = pourcentage_mois_ecoule - pourcentage_tests_effectues
    
    # Déterminer la couleur de l'indicateur
    # Vert : en avance ou léger retard (< 15% de retard)
    # Jaune : retard modéré (15-35% de retard)
    # Rouge : gros retard (> 35% de retard)
    if retard < 15:
        indicateur_couleur = "green"
        indicateur_statut = "Dans les temps"
    elif retard < 35:
        indicateur_couleur = "yellow"
        indicateur_statut = "Retard modéré"
    else:
        indicateur_couleur = "red"
        indicateur_statut = "Retard important"
    
    # Calculer la moyenne de tests réalisés par jour
    jour_actuel = now.day
    moyenne_tests_par_jour = tests_effectues / jour_actuel if jour_actuel > 0 else 0
    
    return {
        "total_programmes": total_programmes,
        "total_partenaires": total_partenaires,
        "total_incidents_ouverts": total_incidents_ouverts,
        "taux_reussite_ts": round(taux_reussite_ts, 2),
        "taux_reussite_tl": round(taux_reussite_tl, 2),
        "tests_manquants": tests_manquants,
        "tests_manquants_count": len(tests_manquants),  # Nombre de combinaisons avec tests manquants
        "tests_manquants_reel": tests_manquants_reel,  # Vrai nombre de tests manquants
        "partenaires_manquants": partenaires_manquants,
        "tests_manquants_j5": tests_manquants_j5,
        "is_j5_alert": is_j5_alert,
        "days_until_end": days_until_end,
        "current_month": month,
        "current_year": year,
        "tests_attendus": tests_attendus,
        "tests_effectues": tests_effectues,
        "pourcentage_mois_ecoule": round(pourcentage_mois_ecoule, 2),
        "pourcentage_tests_effectues": round(pourcentage_tests_effectues, 2),
        "retard": round(retard, 2),
        "indicateur_couleur": indicateur_couleur,
        "indicateur_statut": indicateur_statut,
        "moyenne_tests_par_jour": round(moyenne_tests_par_jour, 2)
    }

# Routes - Export Bilan Partenaire
@api_router.get("/export/bilan-partenaire")
async def export_bilan_partenaire(
    partenaire_id: str = Query(...),
    date_debut: str = Query(...),
    date_fin: str = Query(...)
):
    # Récupérer le partenaire
    partenaire = await db.partenaires.find_one({"id": partenaire_id}, {"_id": 0})
    if not partenaire:
        raise HTTPException(status_code=404, detail="Partenaire non trouvé")
    
    # Récupérer tous les programmes
    programmes = await db.programmes.find({}, {"_id": 0}).to_list(1000)
    programmes_dict = {p['id']: p['nom'] for p in programmes}
    
    # Query pour tous les tests du partenaire dans la période (tous programmes)
    query = {
        'partenaire_id': partenaire_id,
        'date_test': {
            '$gte': date_debut,
            '$lte': date_fin
        }
    }
    
    # Récupérer les tests site et ligne
    tests_site = await db.tests_site.find(query, {"_id": 0}).to_list(10000)
    tests_ligne = await db.tests_ligne.find(query, {"_id": 0}).to_list(10000)
    
    # Create CSV with both test types
    output = io.StringIO()
    
    # Header
    output.write(f"BILAN PARTENAIRE: {partenaire['nom']}\n")
    output.write(f"Période: du {date_debut} au {date_fin}\n")
    output.write(f"Remise minimum attendue: {partenaire.get('remise_minimum', 'Non définie')}%\n")
    output.write("\n")
    
    # Tests Site
    output.write("=== TESTS SITE ===\n")
    if tests_site:
        fieldnames_site = ['Date', 'Programme', 'Application remise', 'Prix public', 'Prix remisé', 
                          '% Remise', 'Naming', 'Cumul codes', 'Commentaire']
        writer = csv.DictWriter(output, fieldnames=fieldnames_site)
        writer.writeheader()
        
        for test in tests_site:
            date_str = test['date_test'] if isinstance(test['date_test'], str) else test['date_test'].isoformat()
            programme_nom = programmes_dict.get(test['programme_id'], test['programme_id'])
            
            writer.writerow({
                'Date': date_str,
                'Programme': programme_nom,
                'Application remise': 'OUI' if test['application_remise'] else 'NON',
                'Prix public': f"{test['prix_public']}€",
                'Prix remisé': f"{test['prix_remise']}€",
                '% Remise': f"{test['pct_remise_calcule']}%",
                'Naming': test.get('naming_constate', ''),
                'Cumul codes': 'OUI' if test['cumul_codes'] else 'NON',
                'Commentaire': test.get('commentaire', '')
            })
    else:
        output.write("Aucun test site sur cette période\n")
    
    output.write("\n\n")
    
    # Tests Ligne
    output.write("=== TESTS LIGNE ===\n")
    if tests_ligne:
        fieldnames_ligne = ['Date', 'Programme', 'Téléphone', 'Messagerie dédiée', 'Décroche dédié',
                           'Délai attente', 'Conseiller', 'Évaluation', 'Offre appliquée', 'Commentaire']
        writer = csv.DictWriter(output, fieldnames=fieldnames_ligne)
        writer.writeheader()
        
        for test in tests_ligne:
            date_str = test['date_test'] if isinstance(test['date_test'], str) else test['date_test'].isoformat()
            programme_nom = programmes_dict.get(test['programme_id'], test['programme_id'])
            
            writer.writerow({
                'Date': date_str,
                'Programme': programme_nom,
                'Téléphone': test['numero_telephone'],
                'Messagerie dédiée': 'OUI' if test['messagerie_vocale_dediee'] else 'NON',
                'Décroche dédié': 'OUI' if test['decroche_dedie'] else 'NON',
                'Délai attente': test['delai_attente'],
                'Conseiller': test.get('nom_conseiller', 'NC'),
                'Évaluation': test['evaluation_accueil'],
                'Offre appliquée': 'OUI' if test['application_offre'] else 'NON',
                'Commentaire': test.get('commentaire', '')
            })
    else:
        output.write("Aucun test ligne sur cette période\n")
    
    output.seek(0)
    return StreamingResponse(
        iter([output.getvalue()]),
        media_type="text/csv",
        headers={"Content-Disposition": f"attachment; filename=bilan_{partenaire['nom']}_{date_debut}_{date_fin}.csv"}
    )

# Routes - Export Bilan Excel Tests Site
@api_router.get("/export/bilan-site-excel")
async def export_bilan_site_excel(
    partenaire_id: str = Query(...),
    date_debut: str = Query(...),
    date_fin: str = Query(...)
):
    from datetime import datetime as dt
    import locale
    
    # Essayer de définir la locale française
    try:
        locale.setlocale(locale.LC_TIME, 'fr_FR.UTF-8')
    except:
        try:
            locale.setlocale(locale.LC_TIME, 'fr_FR')
        except:
            pass  # Si pas de locale française, on utilisera l'anglais
    
    # Récupérer le partenaire
    partenaire = await db.partenaires.find_one({"id": partenaire_id}, {"_id": 0})
    if not partenaire:
        raise HTTPException(status_code=404, detail="Partenaire non trouvé")
    
    # Récupérer tous les programmes
    programmes = await db.programmes.find({}, {"_id": 0}).to_list(1000)
    programmes_dict = {p['id']: p['nom'] for p in programmes}
    
    # Query pour tous les tests site du partenaire dans la période
    query = {
        'partenaire_id': partenaire_id,
        'date_test': {
            '$gte': date_debut,
            '$lte': date_fin
        }
    }
    
    tests_site = await db.tests_site.find(query, {"_id": 0}).to_list(10000)
    
    # Grouper les tests par programme
    tests_par_programme = {}
    for test in tests_site:
        prog_id = test['programme_id']
        if prog_id not in tests_par_programme:
            tests_par_programme[prog_id] = []
        tests_par_programme[prog_id].append(test)
    
    # Créer le workbook Excel
    wb = Workbook()
    wb.remove(wb.active)  # Supprimer la feuille par défaut
    
    # Styles
    header_font = Font(name='Calibri', size=11, bold=True, color='FFFFFF')
    header_fill = PatternFill(start_color='C00000', end_color='C00000', fill_type='solid')
    header_alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
    
    title_font = Font(name='Calibri', size=14, bold=True, color='C00000')
    title_alignment = Alignment(horizontal='center', vertical='center')
    
    cell_alignment_center = Alignment(horizontal='center', vertical='center')
    cell_alignment_left = Alignment(horizontal='left', vertical='center')
    
    border_style = Border(
        left=Side(style='thin', color='000000'),
        right=Side(style='thin', color='000000'),
        top=Side(style='thin', color='000000'),
        bottom=Side(style='thin', color='000000')
    )
    
    # Mapper mois en français
    mois_fr = {
        1: 'Janvier', 2: 'Février', 3: 'Mars', 4: 'Avril',
        5: 'Mai', 6: 'Juin', 7: 'Juillet', 8: 'Août',
        9: 'Septembre', 10: 'Octobre', 11: 'Novembre', 12: 'Décembre'
    }
    
    # Créer une feuille par programme
    for prog_id, tests in tests_par_programme.items():
        programme_nom = programmes_dict.get(prog_id, prog_id)
        
        # Créer la feuille (limiter le nom à 31 caractères pour Excel)
        sheet_name = f"{partenaire['nom'][:15]} - {programme_nom[:12]}"
        ws = wb.create_sheet(title=sheet_name)
        
        # Titre principal (ligne 1 fusionnée)
        ws.merge_cells('A1:F1')
        title_cell = ws['A1']
        title_cell.value = f"TESTS SITE – {partenaire['nom']} – {programme_nom}"
        title_cell.font = title_font
        title_cell.alignment = title_alignment
        ws.row_dimensions[1].height = 25
        
        # En-têtes (ligne 2) - NOUVELLES COLONNES
        headers = ['MOIS', 'DATE EXACTE', 'APPLICATION DE LA REMISE', 
                   'Application claire (Prix GP vs Prix remisé)', 'Naming de la remise', 
                   'Cumul des codes promos']
        
        for col_num, header in enumerate(headers, 1):
            cell = ws.cell(row=2, column=col_num)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_alignment
            cell.border = border_style
        
        ws.row_dimensions[2].height = 30
        
        # Données
        for row_num, test in enumerate(tests, 3):
            # Date
            date_str = test['date_test'] if isinstance(test['date_test'], str) else test['date_test'].isoformat()
            try:
                date_obj = dt.fromisoformat(date_str.replace('Z', '+00:00'))
                # MOIS : Février-2025
                mois_nom = mois_fr.get(date_obj.month, date_obj.strftime('%B'))
                mois_formatted = f"{mois_nom}-{date_obj.year}"
                # DATE EXACTE : 15/02/2025
                date_formatted = date_obj.strftime('%d/%m/%Y')
            except:
                mois_formatted = date_str[:7]
                date_formatted = date_str
            
            # Application claire de la remise
            prix_gp = test.get('prix_public', 0)
            prix_remise = test.get('prix_remise', 0)
            application_claire = f"{prix_gp}€ vs {prix_remise}€"
            
            row_data = [
                mois_formatted,  # MOIS
                date_formatted,  # DATE EXACTE
                'Oui' if test.get('application_remise') else 'Non',  # APPLICATION DE LA REMISE
                application_claire,  # Application claire
                test.get('naming_constate', ''),  # Naming de la remise
                'Oui' if test.get('cumul_codes') else 'Non',  # Cumul des codes promos
            ]
            
            for col_num, value in enumerate(row_data, 1):
                cell = ws.cell(row=row_num, column=col_num)
                cell.value = value
                cell.border = border_style
                cell.alignment = cell_alignment_center
        
        # Ajuster largeurs de colonnes
        column_widths = {
            'A': 18,  # MOIS
            'B': 15,  # DATE EXACTE
            'C': 25,  # APPLICATION DE LA REMISE
            'D': 30,  # Application claire
            'E': 30,  # Naming
            'F': 25   # Cumul codes
        }
        
        for col_letter, width in column_widths.items():
            ws.column_dimensions[col_letter].width = width
    
    # Sauvegarder dans un buffer
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    
    # Nom du fichier : export_azureva_[mois-année].xlsx
    try:
        date_obj_debut = dt.fromisoformat(date_debut.replace('Z', '+00:00'))
        mois_nom = mois_fr.get(date_obj_debut.month, date_obj_debut.strftime('%B')).lower()
        filename = f"export_{partenaire['nom'].lower().replace(' ', '_')}_{mois_nom}-{date_obj_debut.year}.xlsx"
    except:
        filename = f"export_{partenaire['nom'].lower().replace(' ', '_')}.xlsx"
    
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

# Routes - Export Bilan Excel Tests Ligne
@api_router.get("/export/bilan-ligne-excel")
async def export_bilan_ligne_excel(
    partenaire_id: str = Query(...),
    date_debut: str = Query(...),
    date_fin: str = Query(...)
):
    from datetime import datetime as dt
    import locale
    
    # Essayer de définir la locale française
    try:
        locale.setlocale(locale.LC_TIME, 'fr_FR.UTF-8')
    except:
        try:
            locale.setlocale(locale.LC_TIME, 'fr_FR')
        except:
            pass
    
    # Récupérer le partenaire
    partenaire = await db.partenaires.find_one({"id": partenaire_id}, {"_id": 0})
    if not partenaire:
        raise HTTPException(status_code=404, detail="Partenaire non trouvé")
    
    # Récupérer tous les programmes
    programmes = await db.programmes.find({}, {"_id": 0}).to_list(1000)
    programmes_dict = {p['id']: p['nom'] for p in programmes}
    
    # Query pour tous les tests ligne du partenaire dans la période
    query = {
        'partenaire_id': partenaire_id,
        'date_test': {
            '$gte': date_debut,
            '$lte': date_fin
        }
    }
    
    tests_ligne = await db.tests_ligne.find(query, {"_id": 0}).to_list(10000)
    
    # Grouper les tests par programme
    tests_par_programme = {}
    for test in tests_ligne:
        prog_id = test['programme_id']
        if prog_id not in tests_par_programme:
            tests_par_programme[prog_id] = []
        tests_par_programme[prog_id].append(test)
    
    # Créer le workbook Excel
    wb = Workbook()
    wb.remove(wb.active)  # Supprimer la feuille par défaut
    
    # Styles
    header_font = Font(name='Calibri', size=11, bold=True, color='FFFFFF')
    header_fill = PatternFill(start_color='C00000', end_color='C00000', fill_type='solid')
    header_alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
    
    title_font = Font(name='Calibri', size=14, bold=True, color='C00000')
    title_alignment = Alignment(horizontal='center', vertical='center')
    
    cell_alignment_center = Alignment(horizontal='center', vertical='center')
    cell_alignment_left = Alignment(horizontal='left', vertical='center')
    
    border_style = Border(
        left=Side(style='thin', color='000000'),
        right=Side(style='thin', color='000000'),
        top=Side(style='thin', color='000000'),
        bottom=Side(style='thin', color='000000')
    )
    
    # Mapper mois en français
    mois_fr = {
        1: 'Janvier', 2: 'Février', 3: 'Mars', 4: 'Avril',
        5: 'Mai', 6: 'Juin', 7: 'Juillet', 8: 'Août',
        9: 'Septembre', 10: 'Octobre', 11: 'Novembre', 12: 'Décembre'
    }
    
    # Créer une feuille par programme
    for prog_id, tests in tests_par_programme.items():
        programme_nom = programmes_dict.get(prog_id, prog_id)
        
        # Créer la feuille (limiter le nom à 31 caractères pour Excel)
        sheet_name = f"{partenaire['nom'][:15]} - {programme_nom[:12]}"
        ws = wb.create_sheet(title=sheet_name)
        
        # Titre principal (ligne 1 fusionnée)
        ws.merge_cells('A1:I1')
        title_cell = ws['A1']
        title_cell.value = f"TESTS LIGNE – {partenaire['nom']} – {programme_nom}"
        title_cell.font = title_font
        title_cell.alignment = title_alignment
        ws.row_dimensions[1].height = 25
        
        # En-têtes (ligne 2) - NOUVELLES COLONNES
        headers = ['MOIS', 'DATE EXACTE', 'Numéro de téléphone', 'Messagerie Vocale dédiée', 
                   'Délai d\'attente', 'Nom du conseiller', 'Décroche dédiée', 
                   'Évaluation de l\'accueil', 'Application de l\'offre']
        
        for col_num, header in enumerate(headers, 1):
            cell = ws.cell(row=2, column=col_num)
            cell.value = header
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_alignment
            cell.border = border_style
        
        ws.row_dimensions[2].height = 30
        
        # Données
        for row_num, test in enumerate(tests, 3):
            # Date
            date_str = test['date_test'] if isinstance(test['date_test'], str) else test['date_test'].isoformat()
            try:
                date_obj = dt.fromisoformat(date_str.replace('Z', '+00:00'))
                # MOIS : Février-2025
                mois_nom = mois_fr.get(date_obj.month, date_obj.strftime('%B'))
                mois_formatted = f"{mois_nom}-{date_obj.year}"
                # DATE EXACTE : 15/02/2025
                date_formatted = date_obj.strftime('%d/%m/%Y')
            except:
                mois_formatted = date_str[:7]
                date_formatted = date_str
            
            row_data = [
                mois_formatted,  # MOIS
                date_formatted,  # DATE EXACTE
                test.get('numero_telephone', ''),  # Numéro de téléphone
                'Oui' if test.get('messagerie_vocale_dediee') else 'Non',  # Messagerie Vocale dédiée
                test.get('delai_attente', ''),  # Délai d'attente
                test.get('nom_conseiller', ''),  # Nom du conseiller
                'Oui' if test.get('decroche_dedie') else 'Non',  # Décroche dédiée
                test.get('evaluation_accueil', ''),  # Évaluation de l'accueil
                'Oui' if test.get('application_offre') else 'Non',  # Application de l'offre
            ]
            
            for col_num, value in enumerate(row_data, 1):
                cell = ws.cell(row=row_num, column=col_num)
                cell.value = value
                cell.border = border_style
                cell.alignment = cell_alignment_center
        
        # Ajuster largeurs de colonnes
        column_widths = {
            'A': 18,  # MOIS
            'B': 15,  # DATE EXACTE
            'C': 20,  # Numéro de téléphone
            'D': 25,  # Messagerie Vocale dédiée
            'E': 18,  # Délai d'attente
            'F': 25,  # Nom du conseiller
            'G': 18,  # Décroche dédiée
            'H': 25,  # Évaluation de l'accueil
            'I': 25   # Application de l'offre
        }
    
        for col_letter, width in column_widths.items():
            ws.column_dimensions[col_letter].width = width
    
    # Sauvegarder dans un buffer
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    
    # Nom du fichier : export_azureva_[mois-année].xlsx
    try:
        date_obj_debut = dt.fromisoformat(date_debut.replace('Z', '+00:00'))
        mois_nom = mois_fr.get(date_obj_debut.month, date_obj_debut.strftime('%B')).lower()
        filename = f"export_{partenaire['nom'].lower().replace(' ', '_')}_{mois_nom}-{date_obj_debut.year}.xlsx"
    except:
        filename = f"export_{partenaire['nom'].lower().replace(' ', '_')}.xlsx"
    
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

# Routes - Export CSV (legacy)
@api_router.get("/export/tests-site")
async def export_tests_site_csv(
    programme_id: Optional[str] = Query(None),
    partenaire_id: Optional[str] = Query(None)
):
    query = {}
    if programme_id:
        query['programme_id'] = programme_id
    if partenaire_id:
        query['partenaire_id'] = partenaire_id
    
    tests = await db.tests_site.find(query, {"_id": 0}).to_list(10000)
    
    # Create CSV
    output = io.StringIO()
    if tests:
        fieldnames = ['id', 'programme_id', 'partenaire_id', 'date_test', 'application_remise',
                     'prix_public', 'prix_remise', 'pct_remise_calcule', 'naming_constate',
                     'cumul_codes', 'commentaire']
        writer = csv.DictWriter(output, fieldnames=fieldnames)
        writer.writeheader()
        for test in tests:
            # Convert datetime to string
            if 'date_test' in test and not isinstance(test['date_test'], str):
                test['date_test'] = test['date_test'].isoformat() if hasattr(test['date_test'], 'isoformat') else str(test['date_test'])
            # Remove created_at from export
            test.pop('created_at', None)
            writer.writerow({k: test.get(k, '') for k in fieldnames})
    
    output.seek(0)
    return StreamingResponse(
        iter([output.getvalue()]),
        media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=tests_site.csv"}
    )

@api_router.get("/export/tests-ligne")
async def export_tests_ligne_csv(
    programme_id: Optional[str] = Query(None),
    partenaire_id: Optional[str] = Query(None)
):
    query = {}
    if programme_id:
        query['programme_id'] = programme_id
    if partenaire_id:
        query['partenaire_id'] = partenaire_id
    
    tests = await db.tests_ligne.find(query, {"_id": 0}).to_list(10000)
    
    # Create CSV
    output = io.StringIO()
    if tests:
        fieldnames = ['id', 'programme_id', 'partenaire_id', 'date_test', 'numero_telephone',
                     'messagerie_vocale_dediee', 'decroche_dedie', 'delai_attente', 'nom_conseiller',
                     'evaluation_accueil', 'application_offre', 'commentaire']
        writer = csv.DictWriter(output, fieldnames=fieldnames)
        writer.writeheader()
        for test in tests:
            # Convert datetime to string
            if 'date_test' in test and not isinstance(test['date_test'], str):
                test['date_test'] = test['date_test'].isoformat() if hasattr(test['date_test'], 'isoformat') else str(test['date_test'])
            # Remove created_at from export
            test.pop('created_at', None)
            writer.writerow({k: test.get(k, '') for k in fieldnames})
    
    output.seek(0)
    return StreamingResponse(
        iter([output.getvalue()]),
        media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=tests_ligne.csv"}
    )

# Routes - Email Templates
@api_router.get("/email-templates", response_model=List[EmailTemplate])
async def get_email_templates():
    templates = await db.email_templates.find().to_list(length=None)
    return templates

@api_router.post("/email-templates", response_model=EmailTemplate)
async def create_email_template(input: EmailTemplateCreate):
    # If this is set as default, unset other defaults
    if input.is_default:
        await db.email_templates.update_many(
            {"is_default": True},
            {"$set": {"is_default": False}}
        )
    
    template = EmailTemplate(**input.model_dump())
    doc = template.model_dump()
    doc['created_at'] = doc['created_at'].isoformat()
    await db.email_templates.insert_one(doc)
    return template

@api_router.put("/email-templates/{template_id}", response_model=EmailTemplate)
async def update_email_template(template_id: str, input: EmailTemplateCreate):
    existing = await db.email_templates.find_one({"id": template_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Template not found")
    
    # If this is set as default, unset other defaults
    if input.is_default:
        await db.email_templates.update_many(
            {"is_default": True, "id": {"$ne": template_id}},
            {"$set": {"is_default": False}}
        )
    
    update_data = input.model_dump()
    await db.email_templates.update_one(
        {"id": template_id},
        {"$set": update_data}
    )
    
    updated = await db.email_templates.find_one({"id": template_id})
    return EmailTemplate(**updated)

@api_router.delete("/email-templates/{template_id}")
async def delete_email_template(template_id: str):
    result = await db.email_templates.delete_one({"id": template_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Template not found")
    return {"message": "Template deleted"}

@api_router.put("/email-templates/{template_id}/set-default")
async def set_default_template(template_id: str):
    # Unset all defaults
    await db.email_templates.update_many(
        {"is_default": True},
        {"$set": {"is_default": False}}
    )
    
    # Set this one as default
    result = await db.email_templates.update_one(
        {"id": template_id},
        {"$set": {"is_default": True}}
    )
    
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Template not found")
    
    return {"message": "Default template set"}

# Routes - User Signatures
@api_router.get("/signatures", response_model=List[UserSignature])
async def get_signatures():
    signatures = await db.signatures.find().to_list(length=None)
    return signatures

@api_router.post("/signatures", response_model=UserSignature)
async def create_signature(input: UserSignatureCreate):
    signature = UserSignature(**input.model_dump())
    doc = signature.model_dump()
    doc['created_at'] = doc['created_at'].isoformat()
    await db.signatures.insert_one(doc)
    return signature

@api_router.put("/signatures/{signature_id}", response_model=UserSignature)
async def update_signature(signature_id: str, input: UserSignatureCreate):
    existing = await db.signatures.find_one({"id": signature_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Signature not found")
    
    update_data = input.model_dump()
    await db.signatures.update_one(
        {"id": signature_id},
        {"$set": update_data}
    )
    
    updated = await db.signatures.find_one({"id": signature_id})
    return UserSignature(**updated)

@api_router.delete("/signatures/{signature_id}")
async def delete_signature(signature_id: str):
    result = await db.signatures.delete_one({"id": signature_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Signature not found")
    return {"message": "Signature deleted"}

# Routes - Email Drafts
@api_router.get("/email-drafts", response_model=List[EmailDraft])
async def get_email_drafts(status: Optional[str] = None):
    query = {}
    if status:
        query['status'] = status
    drafts = await db.email_drafts.find(query).to_list(length=None)
    return drafts

@api_router.get("/email-drafts/{draft_id}", response_model=EmailDraft)
async def get_email_draft(draft_id: str):
    draft = await db.email_drafts.find_one({"id": draft_id})
    if not draft:
        raise HTTPException(status_code=404, detail="Draft not found")
    return EmailDraft(**draft)

@api_router.post("/email-drafts", response_model=EmailDraft)
async def create_email_draft(input: EmailDraftCreate):
    draft = EmailDraft(**input.model_dump())
    doc = draft.model_dump()
    doc['created_at'] = doc['created_at'].isoformat()
    await db.email_drafts.insert_one(doc)
    return draft

@api_router.put("/email-drafts/{draft_id}", response_model=EmailDraft)
async def update_email_draft(draft_id: str, input: EmailDraftCreate):
    existing = await db.email_drafts.find_one({"id": draft_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Draft not found")
    
    update_data = input.model_dump()
    await db.email_drafts.update_one(
        {"id": draft_id},
        {"$set": update_data}
    )
    
    updated = await db.email_drafts.find_one({"id": draft_id})
    return EmailDraft(**updated)

@api_router.delete("/email-drafts/{draft_id}")
async def delete_email_draft(draft_id: str):
    result = await db.email_drafts.delete_one({"id": draft_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Draft not found")
    return {"message": "Draft deleted"}

class SendEmailRequest(BaseModel):
    signature_id: Optional[str] = None

@api_router.post("/email-drafts/{draft_id}/send")
async def send_email_draft(draft_id: str, request: SendEmailRequest):
    # Get draft
    draft = await db.email_drafts.find_one({"id": draft_id})
    if not draft:
        raise HTTPException(status_code=404, detail="Draft not found")
    
    if draft['status'] == 'sent':
        raise HTTPException(status_code=400, detail="Draft already sent")
    
    # Get signature if provided
    signature_text = ""
    if request.signature_id:
        signature = await db.signatures.find_one({"id": request.signature_id})
        if signature:
            signature_text = signature['signature_text']
    
    # Send email
    result = await send_email_smtp(
        draft['recipient'],
        draft['subject'],
        draft['body'],
        signature_text
    )
    
    if result['status'] == 'success':
        # Update draft status
        await db.email_drafts.update_one(
            {"id": draft_id},
            {"$set": {
                "status": "sent",
                "sent_at": datetime.now(timezone.utc).isoformat()
            }}
        )
        
        # Create history entry
        history = EmailHistory(
            alerte_id=draft['alerte_id'],
            draft_id=draft_id,
            recipient=draft['recipient'],
            subject=draft['subject'],
            body=draft['body'] + f"\n\n{signature_text}" if signature_text else draft['body'],
            status='success'
        )
        history_doc = history.model_dump()
        history_doc['sent_at'] = history_doc['sent_at'].isoformat()
        await db.email_history.insert_one(history_doc)
        
        # Update alerte status to indicate contact was made
        await db.alertes.update_one(
            {"id": draft['alerte_id']},
            {"$set": {"statut": "resolu"}}
        )
        
        return {"message": "Email sent successfully", "status": "success"}
    else:
        # Log failure in history
        history = EmailHistory(
            alerte_id=draft['alerte_id'],
            draft_id=draft_id,
            recipient=draft['recipient'],
            subject=draft['subject'],
            body=draft['body'],
            status='failed',
            error_message=result['message']
        )
        history_doc = history.model_dump()
        history_doc['sent_at'] = history_doc['sent_at'].isoformat()
        await db.email_history.insert_one(history_doc)
        
        raise HTTPException(status_code=500, detail=result['message'])

# Routes - Email History
@api_router.get("/email-history", response_model=List[EmailHistory])
async def get_email_history(alerte_id: Optional[str] = None):
    query = {}
    if alerte_id:
        query['alerte_id'] = alerte_id
    history = await db.email_history.find(query).sort("sent_at", -1).to_list(length=None)
    return history

# =====================
# Authentication Routes
# =====================

@api_router.post("/auth/login", response_model=Token)
async def login(login_request: LoginRequest):
    """Authenticate user and return JWT token"""
    user = await db.users.find_one({"email": login_request.email})
    if not user:
        raise HTTPException(status_code=401, detail="Email ou mot de passe incorrect")
    
    if not verify_password(login_request.password, user['password_hash']):
        raise HTTPException(status_code=401, detail="Email ou mot de passe incorrect")
    
    if not user.get('is_active', True):
        raise HTTPException(status_code=400, detail="Compte utilisateur désactivé")
    
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={"sub": user['email']}, expires_delta=access_token_expires
    )
    return {"access_token": access_token, "token_type": "bearer"}

@api_router.post("/auth/register", response_model=User)
async def register(user_create: UserCreate, current_user: User = Depends(get_current_active_user)):
    """Register a new user (Admin only)"""
    if current_user.role != UserRole.admin:
        raise HTTPException(status_code=403, detail="Seuls les administrateurs peuvent créer des utilisateurs")
    
    # Check if user already exists
    existing_user = await db.users.find_one({"email": user_create.email})
    if existing_user:
        raise HTTPException(status_code=400, detail="Cet email est déjà utilisé")
    
    # Create new user
    user = User(
        email=user_create.email,
        nom=user_create.nom,
        prenom=user_create.prenom,
        role=user_create.role,
        is_active=user_create.is_active,
        programme_id=user_create.programme_id,
        partenaire_id=user_create.partenaire_id,
        programme_ids=user_create.programme_ids
    )
    
    user_dict = user.model_dump()
    user_dict['password_hash'] = get_password_hash(user_create.password)
    user_dict['created_at'] = user_dict['created_at'].isoformat()
    
    await db.users.insert_one(user_dict)
    return user

@api_router.post("/auth/init-admin")
async def init_admin():
    """Initialize default admin user (only if no users exist)"""
    count = await db.users.count_documents({})
    if count > 0:
        raise HTTPException(status_code=400, detail="Des utilisateurs existent déjà")
    
    admin = User(
        email="admin@hubblindtests.com",
        nom="Admin",
        prenom="Super",
        role=UserRole.admin,
        is_active=True
    )
    
    admin_dict = admin.model_dump()
    admin_dict['password_hash'] = get_password_hash("admin123")
    admin_dict['created_at'] = admin_dict['created_at'].isoformat()
    
    await db.users.insert_one(admin_dict)
    return {"message": "Administrateur créé avec succès", "email": "admin@hubblindtests.com", "password": "admin123"}

# =====================
# User Management Routes
# =====================

@api_router.get("/users/me", response_model=User)
async def get_my_profile(current_user: User = Depends(get_current_active_user)):
    """Get current user profile"""
    return current_user

@api_router.get("/users", response_model=List[User])
async def get_users(current_user: User = Depends(get_current_active_user)):
    """Get all users (Admin only)"""
    if current_user.role != UserRole.admin:
        raise HTTPException(status_code=403, detail="Accès réservé aux administrateurs")
    
    users = await db.users.find().to_list(length=None)
    return [User(**user) for user in users]

@api_router.get("/users/{user_id}", response_model=User)
async def get_user(user_id: str, current_user: User = Depends(get_current_active_user)):
    """Get user by ID (Admin only)"""
    if current_user.role != UserRole.admin:
        raise HTTPException(status_code=403, detail="Accès réservé aux administrateurs")
    
    user = await db.users.find_one({"id": user_id})
    if not user:
        raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
    return User(**user)

@api_router.post("/users", response_model=User)
async def create_user(user_create: UserCreate, current_user: User = Depends(get_current_active_user)):
    """Create a new user (Admin only)"""
    if current_user.role != UserRole.admin:
        raise HTTPException(status_code=403, detail="Seuls les administrateurs peuvent créer des utilisateurs")
    
    # Check if user already exists
    existing_user = await db.users.find_one({"email": user_create.email})
    if existing_user:
        raise HTTPException(status_code=400, detail="Cet email est déjà utilisé")
    
    # Create new user
    user = User(
        email=user_create.email,
        nom=user_create.nom,
        prenom=user_create.prenom,
        role=user_create.role,
        is_active=user_create.is_active,
        programme_id=user_create.programme_id,
        partenaire_id=user_create.partenaire_id,
        programme_ids=user_create.programme_ids
    )
    
    user_dict = user.model_dump()
    user_dict['password_hash'] = get_password_hash(user_create.password)
    user_dict['created_at'] = user_dict['created_at'].isoformat()
    
    await db.users.insert_one(user_dict)
    return user

@api_router.put("/users/{user_id}", response_model=User)
async def update_user(user_id: str, user_update: UserUpdate, current_user: User = Depends(get_current_active_user)):
    """Update user (Admin only)"""
    if current_user.role != UserRole.admin:
        raise HTTPException(status_code=403, detail="Seuls les administrateurs peuvent modifier des utilisateurs")
    
    user = await db.users.find_one({"id": user_id})
    if not user:
        raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
    
    update_data = {}
    if user_update.nom is not None:
        update_data['nom'] = user_update.nom
    if user_update.prenom is not None:
        update_data['prenom'] = user_update.prenom
    if user_update.role is not None:
        update_data['role'] = user_update.role
    if user_update.is_active is not None:
        update_data['is_active'] = user_update.is_active
    if user_update.password is not None:
        update_data['password_hash'] = get_password_hash(user_update.password)
    
    if update_data:
        await db.users.update_one({"id": user_id}, {"$set": update_data})
    
    updated_user = await db.users.find_one({"id": user_id})
    return User(**updated_user)

@api_router.delete("/users/{user_id}")
async def delete_user(user_id: str, current_user: User = Depends(get_current_active_user)):
    """Delete user (Admin only)"""
    if current_user.role != UserRole.admin:
        raise HTTPException(status_code=403, detail="Seuls les administrateurs peuvent supprimer des utilisateurs")
    
    if user_id == current_user.id:
        raise HTTPException(status_code=400, detail="Vous ne pouvez pas supprimer votre propre compte")
    
    result = await db.users.delete_one({"id": user_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
    
    return {"message": "Utilisateur supprimé avec succès"}

@api_router.get("/users/stats/all")
async def get_users_stats(current_user: User = Depends(get_current_active_user)):
    """Get statistics for all users"""
    if current_user.role != UserRole.admin:
        raise HTTPException(status_code=403, detail="Accès réservé aux administrateurs")
    
    users = await db.users.find().to_list(length=None)
    stats = []
    
    for user in users:
        user_id = user['id']
        
        # Count tests created by user
        tests_site = await db.tests_site.count_documents({"user_id": user_id})
        tests_ligne = await db.tests_ligne.count_documents({"user_id": user_id})
        total_tests = tests_site + tests_ligne
        
        # Count alertes handled
        alertes = await db.alertes.count_documents({"user_id": user_id})
        
        stats.append({
            "user": User(**user),
            "tests_site_count": tests_site,
            "tests_ligne_count": tests_ligne,
            "total_tests": total_tests,
            "incidents_count": alertes
        })
    
    return stats

# Helper functions for PPT generation
def clone_slide_xml(prs_zip_path, slide_index, output_path):
    """Clone a slide by manipulating the PPTX ZIP structure"""
    # Extract the pptx
    with zipfile.ZipFile(prs_zip_path, 'r') as zip_ref:
        zip_ref.extractall(output_path)
    
    # Read slide XML
    slide_path = f'{output_path}/ppt/slides/slide{slide_index + 1}.xml'
    with open(slide_path, 'rb') as f:
        slide_tree = etree.parse(f)
    
    return slide_tree

def replace_text_in_shape(shape, replacements):
    """Replace placeholder text in a shape"""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                for key, value in replacements.items():
                    if key in original_text:
                        run.text = original_text.replace(key, str(value))
                        original_text = run.text

def replace_text_in_slide(slide, replacements):
    """Replace all placeholder text in a slide"""
    for shape in slide.shapes:
        replace_text_in_shape(shape, replacements)
        # Also check tables
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            original_text = run.text
                            for key, value in replacements.items():
                                if key in original_text:
                                    run.text = original_text.replace(key, str(value))
                                    original_text = run.text

def fill_table_with_data(table, data_rows, header_rows=1):
    """Fill a table with data rows - simplified approach"""
    # If no data, return without modifying
    if not data_rows:
        logging.info("No data to fill in table")
        return
    
    # Add data rows directly without clearing
    for row_data in data_rows:
        try:
            row = table.rows.add()
            for i, cell_value in enumerate(row_data):
                if i < len(row.cells):
                    row.cells[i].text = str(cell_value)
        except Exception as e:
            logging.error(f"Error adding row to table: {str(e)}")
            continue

def format_french_month(date_obj):
    """Format date to French month name"""
    months_fr = {
        1: 'Janvier', 2: 'Février', 3: 'Mars', 4: 'Avril',
        5: 'Mai', 6: 'Juin', 7: 'Juillet', 8: 'Août',
        9: 'Septembre', 10: 'Octobre', 11: 'Novembre', 12: 'Décembre'
    }
    return months_fr.get(date_obj.month, str(date_obj.month))

# Routes - Authentication
# Routes - DEBUG Bilan Partenaire
@api_router.get("/debug/bilan-partenaire-analysis")
async def debug_bilan_partenaire_analysis(
    partenaire_id: str = Query(...),
    period_type: str = Query(default="month"),
    year: int = Query(default=2025),
    month: int = Query(default=10)
):
    """Ultra-verbose debug analysis for PowerPoint generation"""
    try:
        debug_report = {
            "SECTION_A_SCHEMA": {},
            "SECTION_B_MAPPING": {},
            "SECTION_C_SAMPLE": {},
            "SECTION_D_PLACEHOLDERS": [],
            "SECTION_E_TABLES_EXEMPLE": [],
            "SECTION_F_DRYRUN_MAP": {},
            "SECTION_G_TABLES_COUNT": {},
            "SECTION_H_ASSERTS": {},
            "SECTION_I_FALLBACK": {}
        }
        
        # PHASE 1 - DATA CONTRACT
        # Get partenaire
        partenaire = await db.partenaires.find_one({"id": partenaire_id})
        if not partenaire:
            return {"error": "Partenaire not found"}
        
        # Get programmes
        programme_ids = partenaire.get('programmes_ids', [])
        programmes = await db.programmes.find({"id": {"$in": programme_ids}}).to_list(length=None)
        programmes = sorted(programmes, key=lambda p: p['nom'])
        
        if not programmes:
            return {"error": "No programmes found"}
        
        programme = programmes[0]
        
        # Calculate period
        today = datetime.now(timezone.utc)
        if period_type == "month":
            date_debut = datetime(year, month, 1, tzinfo=timezone.utc)
            if month == 12:
                date_fin = datetime(year + 1, 1, 1, tzinfo=timezone.utc)
            else:
                date_fin = datetime(year, month + 1, 1, tzinfo=timezone.utc)
            period_label = f"{format_french_month(date_debut)} {year}"
        elif period_type == "year":
            date_debut = datetime(year, 1, 1, tzinfo=timezone.utc)
            date_fin = datetime(year + 1, 1, 1, tzinfo=timezone.utc)
            period_label = f"Année {year}"
        else:
            date_fin = today
            date_debut = datetime(today.year - 1, today.month, 1, tzinfo=timezone.utc)
            period_label = f"Année glissante"
        
        # Get tests
        tests_site = await db.tests_site.find({
            "programme_id": programme['id'],
            "partenaire_id": partenaire_id,
            "date_test": {"$gte": date_debut.isoformat(), "$lt": date_fin.isoformat()}
        }).sort("date_test", 1).to_list(length=3)
        
        tests_ligne = await db.tests_ligne.find({
            "programme_id": programme['id'],
            "partenaire_id": partenaire_id,
            "date_test": {"$gte": date_debut.isoformat(), "$lt": date_fin.isoformat()}
        }).sort("date_test", 1).to_list(length=3)
        
        # SECTION A - SCHEMA
        debug_report["SECTION_A_SCHEMA"] = {
            "partner": {
                "type": "object",
                "keys": list(partenaire.keys()) if partenaire else [],
                "sample_id": partenaire.get('id', 'N/A'),
                "sample_name": partenaire.get('nom', 'N/A')
            },
            "program": {
                "type": "object",
                "keys": list(programme.keys()) if programme else [],
                "sample_id": programme.get('id', 'N/A'),
                "sample_name": programme.get('nom', 'N/A')
            },
            "period": {
                "type": "object",
                "year": year,
                "month": month if period_type == "month" else None,
                "label": period_label,
                "date_debut": date_debut.isoformat(),
                "date_fin": date_fin.isoformat()
            },
            "testsSites": {
                "type": "array",
                "count": len(tests_site),
                "keys": list(tests_site[0].keys()) if tests_site else []
            },
            "testsLignes": {
                "type": "array",
                "count": len(tests_ligne),
                "keys": list(tests_ligne[0].keys()) if tests_ligne else []
            }
        }
        
        # SECTION B - MAPPING
        debug_report["SECTION_B_MAPPING"] = {
            "{PartnerName}": f"partner.nom = '{partenaire.get('nom', 'N/A')}'",
            "{ProgramName}": f"program.nom = '{programme.get('nom', 'N/A')}'",
            "{Mois du trigger + année du trigger}": f"period_label = '{period_label}'",
            "{moyenne des tests sites réussis}": "calculated from testsSites.application_remise",
            "{moyenne des tests lignes réussis}": "calculated from testsLignes.application_offre",
            "{temps d'attente/nombre de test effectués}": "calculated from testsLignes.delai_attente",
            "Bilan du": f"current_date = '{datetime.now(timezone.utc).strftime('%d/%m/%Y')}'",
            "{nom de la remise}": "NOT MAPPED",
            "{info case MD}": "NOT MAPPED",
            "{info case DD}": "NOT MAPPED",
            "{commentaire sur l'accueil}": "NOT MAPPED"
        }
        
        # SECTION C - SAMPLE DATA
        debug_report["SECTION_C_SAMPLE"] = {
            "partnerName": partenaire.get('nom', 'N/A'),
            "programName": programme.get('nom', 'N/A'),
            "period": period_label,
            "testsSites": [
                {k: str(v) for k, v in test.items() if k != '_id'} 
                for test in tests_site[:3]
            ],
            "testsLignes": [
                {k: str(v) for k, v in test.items() if k != '_id'} 
                for test in tests_ligne[:3]
            ]
        }
        
        # PHASE 2 - SCAN TEMPLATE
        template_path = TEMPLATE_DIR / "Bilan_Blindtest_template.pptx"
        if template_path.exists():
            prs = Presentation(str(template_path))
            
            placeholders_found = []
            tables_found = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_type = "Unknown"
                
                # Check slide title
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text
                        
                        # Determine slide type
                        if 'Sites' in text:
                            slide_type = "Tests Sites"
                        elif 'Ligne' in text:
                            slide_type = "Tests Ligne"
                        elif 'SAV' in text:
                            slide_type = "Rapport SAV"
                        elif 'Blind test' in text and slide_type == "Unknown":
                            slide_type = "Vue d'ensemble"
                        
                        # Find placeholders
                        import re
                        patterns = [r'\{([^}]+)\}', r'\{\{([^}]+)\}\}', r'\[\[([^]]+)\]\]']
                        for pattern in patterns:
                            matches = re.findall(pattern, text)
                            for match in matches:
                                placeholders_found.append({
                                    "slide": slide_idx,
                                    "slideType": slide_type,
                                    "location": "text_frame",
                                    "placeholder": f"{{{match}}}",
                                    "mappedKey": f"partner.nom / program.nom / period" if "Partner" in match or "Program" in match else "CALCULATED" if "moyenne" in match or "temps" in match else "NOT MAPPED",
                                    "status": "OK" if ("Partner" in match or "Program" in match or "trigger" in match or "moyenne" in match or "temps" in match) else "NON MAPPÉ"
                                })
                    
                    # Check for tables
                    if shape.has_table:
                        tables_found.append({
                            "slide": slide_idx,
                            "slideType": slide_type,
                            "rows": len(shape.table.rows),
                            "columns": len(shape.table.rows[0].cells) if shape.table.rows else 0,
                            "recommendation": "KEEP_AND_FILL" if slide_type in ["Tests Sites", "Tests Ligne"] else "KEEP_AS_IS"
                        })
            
            debug_report["SECTION_D_PLACEHOLDERS"] = placeholders_found
            debug_report["SECTION_E_TABLES_EXEMPLE"] = tables_found
        else:
            debug_report["SECTION_D_PLACEHOLDERS"] = ["Template file not found"]
            debug_report["SECTION_E_TABLES_EXEMPLE"] = ["Template file not found"]
        
        # PHASE 3 - DRY RUN
        # Calculate statistics
        total_tests_site = len(tests_site)
        tests_site_reussis = len([t for t in tests_site if t.get('application_remise', False)])
        pct_site = round((tests_site_reussis / total_tests_site * 100), 1) if total_tests_site > 0 else 0
        
        total_tests_ligne = len(tests_ligne)
        tests_ligne_reussis = len([t for t in tests_ligne if t.get('application_offre', False)])
        pct_ligne = round((tests_ligne_reussis / total_tests_ligne * 100), 1) if total_tests_ligne > 0 else 0
        
        delais = []
        for t in tests_ligne:
            if t.get('delai_attente'):
                try:
                    parts = t['delai_attente'].split(':')
                    if len(parts) == 2:
                        minutes = int(parts[0])
                        seconds = int(parts[1])
                        delais.append(minutes * 60 + seconds)
                except:
                    pass
        avg_delai = sum(delais) / len(delais) if delais else 0
        avg_delai_str = f"{int(avg_delai // 60):02d}:{int(avg_delai % 60):02d}"
        
        # SECTION F - DRYRUN MAP
        debug_report["SECTION_F_DRYRUN_MAP"] = {
            "{PartnerName}": partenaire.get('nom', 'EMPTY'),
            "{ProgramName}": programme.get('nom', 'EMPTY'),
            "{Mois du trigger + année du trigger}": period_label,
            "{moyenne des tests sites réussis}": f"{pct_site}%",
            "{moyenne des tests lignes réussis}": f"{pct_ligne}%",
            "{temps d'attente/nombre de test effectués}": avg_delai_str,
            "Bilan du": f"Bilan du {datetime.now(timezone.utc).strftime('%d/%m/%Y')}"
        }
        
        # Prepare table rows
        site_rows = []
        for test in tests_site:
            try:
                test_date = datetime.fromisoformat(test['date_test'])
                pct_remise = test.get('pct_remise_calcule', 0)
                site_rows.append([
                    format_french_month(test_date),
                    test_date.strftime('%d/%m/%Y'),
                    f"{test.get('prix_public', 0):.2f} € VS {test.get('prix_remise', 0):.2f} €",
                    'OUI' if test.get('application_remise') else 'NON',
                    test.get('naming_constate', 'N/A'),
                    'OUI' if test.get('cumul_codes') else 'NON'
                ])
            except:
                pass
        
        ligne_rows = []
        for test in tests_ligne:
            try:
                test_date = datetime.fromisoformat(test['date_test'])
                ligne_rows.append([
                    format_french_month(test_date),
                    test_date.strftime('%d/%m/%Y'),
                    test.get('numero_telephone', 'N/A'),
                    'OUI' if test.get('messagerie_vocale_dediee') else 'NON',
                    test.get('delai_attente', 'N/A'),
                    test.get('nom_conseiller', 'N/A'),
                    'OUI' if test.get('decroche_dedie') else 'NON',
                    test.get('evaluation_accueil', 'N/A'),
                    'OUI' if test.get('application_offre') else 'NON'
                ])
            except:
                pass
        
        # SECTION G - TABLES COUNT
        debug_report["SECTION_G_TABLES_COUNT"] = {
            "sitesRowsPrepared": len(site_rows),
            "lignesRowsPrepared": len(ligne_rows)
        }
        
        # SECTION H - ASSERTIONS
        assertions = {}
        
        # A1: No unmapped placeholders
        unmapped = [p for p in debug_report["SECTION_D_PLACEHOLDERS"] if isinstance(p, dict) and p.get("status") == "NON MAPPÉ"]
        assertions["A1_NO_UNMAPPED"] = {
            "status": "PASS" if len(unmapped) == 0 else "FAIL",
            "detail": f"Found {len(unmapped)} unmapped placeholders" if unmapped else "All placeholders mapped"
        }
        
        # A2: No remaining braces in values
        has_braces = any("{" in str(v) for v in debug_report["SECTION_F_DRYRUN_MAP"].values())
        assertions["A2_NO_BRACES"] = {
            "status": "FAIL" if has_braces else "PASS",
            "detail": "Some values still contain braces" if has_braces else "No braces in resolved values"
        }
        
        # A3/T1/T2: Table data validation
        if len(tests_site) > 0 or len(tests_ligne) > 0:
            assertions["T1_SITES_NOT_EMPTY"] = {
                "status": "PASS" if len(site_rows) > 0 else "FAIL",
                "detail": f"Prepared {len(site_rows)} rows for Sites" if len(site_rows) > 0 else "TABLE_SITES_VIDE_ANORMALE"
            }
            assertions["T2_LIGNES_NOT_EMPTY"] = {
                "status": "PASS" if len(ligne_rows) > 0 else "FAIL",
                "detail": f"Prepared {len(ligne_rows)} rows for Lignes" if len(ligne_rows) > 0 else "TABLE_LIGNES_VIDE_ANORMALE"
            }
        
        # A4: Partner and program names not empty
        assertions["A4_NAMES_NOT_EMPTY"] = {
            "status": "PASS" if partenaire.get('nom') and programme.get('nom') else "FAIL",
            "detail": f"Partner: {partenaire.get('nom', 'EMPTY')}, Program: {programme.get('nom', 'EMPTY')}"
        }
        
        debug_report["SECTION_H_ASSERTS"] = assertions
        
        # SECTION I - FALLBACK
        debug_report["SECTION_I_FALLBACK"] = {
            "F1_TEXT_ANCHORS_NEEDED": "NO",
            "F2_TABLE_RECREATION_NEEDED": "NO",
            "reason": "Current approach adds rows to existing tables. If tables don't exist in template, F2 would be needed.",
            "recommendation": "Verify template contains tables on Sites and Lignes slides"
        }
        
        return debug_report
        
    except Exception as e:
        logging.error(f"Debug analysis error: {str(e)}")
        import traceback
        return {
            "error": str(e),
            "traceback": traceback.format_exc()
        }

# Routes - Bilan Partenaire PPT Export (FROM SCRATCH)
@api_router.get("/export/bilan-partenaire-ppt")
async def export_bilan_partenaire_ppt(
    partenaire_id: str = Query(...),
    date_debut: str = Query(...),
    date_fin: str = Query(...)
):
    """Generate PowerPoint report from scratch with real data"""
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        # === GET DATA ===
        partenaire = await db.partenaires.find_one({"id": partenaire_id})
        if not partenaire:
            raise HTTPException(status_code=404, detail="Partenaire not found")
        
        partner_name = partenaire.get('nom', '')
        
        # Get programmes
        programme_ids = partenaire.get('programmes_ids', [])
        programmes = await db.programmes.find({"id": {"$in": programme_ids}}).to_list(length=None)
        programmes = sorted(programmes, key=lambda p: p['nom'])
        
        if not programmes:
            raise HTTPException(status_code=404, detail="No programmes found")
        
        # Parse dates
        date_debut_obj = datetime.fromisoformat(date_debut).replace(tzinfo=timezone.utc)
        date_fin_obj = datetime.fromisoformat(date_fin).replace(tzinfo=timezone.utc, hour=23, minute=59, second=59)
        
        # Generate period label
        if date_debut_obj.year == date_fin_obj.year and date_debut_obj.month == date_fin_obj.month:
            period_label = f"{format_french_month(date_debut_obj)} {date_debut_obj.year}"
        elif date_debut_obj.year == date_fin_obj.year:
            period_label = f"{format_french_month(date_debut_obj)} - {format_french_month(date_fin_obj)} {date_debut_obj.year}"
        else:
            period_label = f"{format_french_month(date_debut_obj)} {date_debut_obj.year} - {format_french_month(date_fin_obj)} {date_fin_obj.year}"
        
        # === CREATE PRESENTATION FROM SCRATCH ===
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        slide_number = 0
        total_slides = len(programmes) * 3  # 3 slides per programme
        
        # === GENERATE SLIDES FOR EACH PROGRAMME ===
        for programme in programmes:
            program_name = programme.get('nom', '')
            
            # Get tests data for THIS programme
            tests_site = await db.tests_site.find({
                "programme_id": programme['id'],
                "partenaire_id": partenaire_id,
                "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
            }).sort("date_test", -1).to_list(length=None)
            
            tests_ligne = await db.tests_ligne.find({
                "programme_id": programme['id'],
                "partenaire_id": partenaire_id,
                "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
            }).sort("date_test", -1).to_list(length=None)
        
        # === CREATE PRESENTATION FROM SCRATCH ===
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        slide_number = 0
        total_slides = len(programmes) * 3  # 3 slides per programme
        
        # === GENERATE SLIDES FOR EACH PROGRAMME ===
        for programme in programmes:
            program_name = programme.get('nom', '')
            
            # Get tests data for THIS programme
            tests_site = await db.tests_site.find({
                "programme_id": programme['id'],
                "partenaire_id": partenaire_id,
                "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
            }).sort("date_test", -1).to_list(length=None)
            
            tests_ligne = await db.tests_ligne.find({
                "programme_id": programme['id'],
                "partenaire_id": partenaire_id,
                "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
            }).sort("date_test", -1).to_list(length=None)
        
        # === CALCULATE STATISTICS ===
        # Sites
        total_tests_site = len(tests_site)
        tests_site_reussis = len([t for t in tests_site if t.get('application_remise', False)])
        pct_site = round((tests_site_reussis / total_tests_site * 100), 1) if total_tests_site > 0 else 0
        
        # Lignes
        total_tests_ligne = len(tests_ligne)
        tests_ligne_reussis = len([t for t in tests_ligne if t.get('application_offre', False)])
        pct_ligne = round((tests_ligne_reussis / total_tests_ligne * 100), 1) if total_tests_ligne > 0 else 0
        
        # Average waiting time
        delais = []
        for t in tests_ligne:
            if t.get('delai_attente'):
                try:
                    parts = t['delai_attente'].split(':')
                    if len(parts) == 2:
                        delais.append(int(parts[0]) * 60 + int(parts[1]))
                except:
                    pass
        avg_delai = sum(delais) / len(delais) if delais else 0
        avg_delai_str = f"{int(avg_delai // 60):02d}:{int(avg_delai % 60):02d}"
        
        # Average accueil
        accueils = [t.get('evaluation_accueil', '') for t in tests_ligne if t.get('evaluation_accueil')]
        accueil_counts = {}
        for acc in accueils:
            accueil_counts[acc] = accueil_counts.get(acc, 0) + 1
        commentaire_accueil = max(accueil_counts, key=accueil_counts.get) if accueil_counts else "—"
        
        # Messagerie/Decroche stats
        md_count = len([t for t in tests_ligne if t.get('messagerie_vocale_dediee')])
        dd_count = len([t for t in tests_ligne if t.get('decroche_dedie')])
        pct_md = round((md_count / total_tests_ligne * 100), 1) if total_tests_ligne > 0 else 0
        pct_dd = round((dd_count / total_tests_ligne * 100), 1) if total_tests_ligne > 0 else 0
        
        # === CREATE PRESENTATION FROM SCRATCH ===
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        slide_number = 0
        total_slides = len(programmes) * 3  # 3 slides per programme
        
        # === GENERATE SLIDES FOR EACH PROGRAMME ===
        for programme in programmes:
            program_name = programme.get('nom', '')
            
            # Get tests data for THIS programme
            tests_site = await db.tests_site.find({
                "programme_id": programme['id'],
                "partenaire_id": partenaire_id,
                "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
            }).sort("date_test", -1).to_list(length=None)
            
            tests_ligne = await db.tests_ligne.find({
                "programme_id": programme['id'],
                "partenaire_id": partenaire_id,
                "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
            }).sort("date_test", -1).to_list(length=None)
            
            # === CALCULATE STATISTICS FOR THIS PROGRAMME ===
            total_tests_site = len(tests_site)
            tests_site_reussis = len([t for t in tests_site if t.get('application_remise', False)])
            pct_site = round((tests_site_reussis / total_tests_site * 100), 1) if total_tests_site > 0 else 0
            
            total_tests_ligne = len(tests_ligne)
            tests_ligne_reussis = len([t for t in tests_ligne if t.get('application_offre', False)])
            pct_ligne = round((tests_ligne_reussis / total_tests_ligne * 100), 1) if total_tests_ligne > 0 else 0
            
            delais = []
            for t in tests_ligne:
                if t.get('delai_attente'):
                    try:
                        parts = t['delai_attente'].split(':')
                        if len(parts) == 2:
                            delais.append(int(parts[0]) * 60 + int(parts[1]))
                    except:
                        pass
            avg_delai = sum(delais) / len(delais) if delais else 0
            avg_delai_str = f"{int(avg_delai // 60):02d}:{int(avg_delai % 60):02d}"
            
            accueils = [t.get('evaluation_accueil', '') for t in tests_ligne if t.get('evaluation_accueil')]
            accueil_counts = {}
            for acc in accueils:
                accueil_counts[acc] = accueil_counts.get(acc, 0) + 1
            commentaire_accueil = max(accueil_counts, key=accueil_counts.get) if accueil_counts else "—"
            
            # === SLIDE 1: VUE D'ENSEMBLE (FOR THIS PROGRAMME) ===
            slide_number += 1
            slide_layout = prs.slide_layouts[6]
            slide1 = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"Blind test – {partner_name} x {program_name}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 128)
        
        # Subtitle - Period
        subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = period_label
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
        
        # Statistics boxes
        y_pos = 2.5
        
        # Tests Sites
        box1 = slide1.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(3.5), Inches(1.5))
        tf1 = box1.text_frame
        tf1.text = f"Tests Sites\n\n{pct_site}% de réussite\n({tests_site_reussis}/{total_tests_site} tests)"
        for para in tf1.paragraphs:
            para.font.size = Pt(16)
            para.alignment = PP_ALIGN.CENTER
        
        # Tests Lignes
        box2 = slide1.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(3.5), Inches(1.5))
        tf2 = box2.text_frame
        tf2.text = f"Tests Ligne\n\n{pct_ligne}% de réussite\n({tests_ligne_reussis}/{total_tests_ligne} tests)"
        for para in tf2.paragraphs:
            para.font.size = Pt(16)
            para.alignment = PP_ALIGN.CENTER
        
        # Additional stats
        box3 = slide1.shapes.add_textbox(Inches(1), Inches(y_pos + 2), Inches(3.5), Inches(1.5))
        tf3 = box3.text_frame
        tf3.text = f"Temps d'attente moyen\n\n{avg_delai_str}"
        for para in tf3.paragraphs:
            para.font.size = Pt(16)
            para.alignment = PP_ALIGN.CENTER
        
        box4 = slide1.shapes.add_textbox(Inches(5.5), Inches(y_pos + 2), Inches(3.5), Inches(1.5))
        tf4 = box4.text_frame
        tf4.text = f"Accueil\n\n{commentaire_accueil}"
        for para in tf4.paragraphs:
            para.font.size = Pt(16)
            para.alignment = PP_ALIGN.CENTER
            
            # Footer
            footer = slide1.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
            footer.text_frame.text = f"Bilan du {datetime.now(timezone.utc).strftime('%d/%m/%Y')} - Page {slide_number}/{total_slides}"
            footer.text_frame.paragraphs[0].font.size = Pt(10)
            footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
            
            # === SLIDE 2: TESTS SITES (FOR THIS PROGRAMME) ===
            slide_number += 1
            slide2 = prs.slides.add_slide(slide_layout)
        
        # Title
        title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        title2.text_frame.text = f"Tests Sites – {partner_name} x {program_name}"
        title2.text_frame.paragraphs[0].font.size = Pt(24)
        title2.text_frame.paragraphs[0].font.bold = True
        
        # Table
        if tests_site:
            rows = min(len(tests_site) + 1, 15)  # Max 14 tests + header
            cols = 7
            table = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
            
            # Header
            headers = ['Date', 'URL', 'Prix Public', 'Prix Remisé', 'Remise', 'OK?', 'Naming']
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
            
            # Data rows
            for idx, test in enumerate(tests_site[:14]):
                if idx + 1 >= rows:
                    break
                try:
                    test_date = datetime.fromisoformat(test['date_test'])
                    table.cell(idx + 1, 0).text = test_date.strftime('%d/%m/%Y')
                    table.cell(idx + 1, 1).text = test.get('url', 'N/A')[:30]
                    table.cell(idx + 1, 2).text = f"{test.get('prix_public', 0):.2f} €"
                    table.cell(idx + 1, 3).text = f"{test.get('prix_remise', 0):.2f} €"
                    table.cell(idx + 1, 4).text = f"{test.get('pct_remise_calcule', 0):.1f}%"
                    table.cell(idx + 1, 5).text = '✓' if test.get('application_remise') else '✗'
                    table.cell(idx + 1, 6).text = test.get('naming_constate', '')[:20]
                    
                    # Font size
                    for col in range(cols):
                        table.cell(idx + 1, col).text_frame.paragraphs[0].font.size = Pt(9)
                except Exception as e:
                    logging.error(f"Error adding test site row: {str(e)}")
        else:
            no_data = slide2.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            no_data.text_frame.text = "Aucun test site disponible pour cette période"
            no_data.text_frame.paragraphs[0].font.size = Pt(18)
            no_data.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Footer
            footer2 = slide2.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
            footer2.text_frame.text = f"Bilan du {datetime.now(timezone.utc).strftime('%d/%m/%Y')} - Page {slide_number}/{total_slides}"
            footer2.text_frame.paragraphs[0].font.size = Pt(10)
            footer2.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
            
            # === SLIDE 3: TESTS LIGNE (FOR THIS PROGRAMME) ===
            slide_number += 1
            slide3 = prs.slides.add_slide(slide_layout)
            
            # Title
            title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
            title3.text_frame.text = f"Tests Ligne – {partner_name} x {program_name}"
            title3.text_frame.paragraphs[0].font.size = Pt(24)
            title3.text_frame.paragraphs[0].font.bold = True
            
            # Table
            if tests_ligne:
                rows = min(len(tests_ligne) + 1, 15)
                cols = 7
                table = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
                
                # Header
                headers = ['Date', 'Téléphone', 'Délai', 'Msg. Vocale', 'Décroche', 'Accueil', 'OK?']
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(11)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
                
                # Data rows
                for idx, test in enumerate(tests_ligne[:14]):
                    if idx + 1 >= rows:
                        break
                    try:
                        test_date = datetime.fromisoformat(test['date_test'])
                        table.cell(idx + 1, 0).text = test_date.strftime('%d/%m/%Y')
                        table.cell(idx + 1, 1).text = test.get('numero_telephone', 'N/A')[:15]
                        table.cell(idx + 1, 2).text = test.get('delai_attente', 'N/A')
                        table.cell(idx + 1, 3).text = '✓' if test.get('messagerie_vocale_dediee') else '✗'
                        table.cell(idx + 1, 4).text = '✓' if test.get('decroche_dedie') else '✗'
                        table.cell(idx + 1, 5).text = test.get('evaluation_accueil', 'N/A')[:15]
                        table.cell(idx + 1, 6).text = '✓' if test.get('application_offre') else '✗'
                        
                        # Font size
                        for col in range(cols):
                            table.cell(idx + 1, col).text_frame.paragraphs[0].font.size = Pt(9)
                    except Exception as e:
                        logging.error(f"Error adding test ligne row: {str(e)}")
            else:
                no_data = slide3.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
                no_data.text_frame.text = "Aucun test ligne disponible pour cette période"
                no_data.text_frame.paragraphs[0].font.size = Pt(18)
                no_data.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Footer
            footer3 = slide3.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
            footer3.text_frame.text = f"Bilan du {datetime.now(timezone.utc).strftime('%d/%m/%Y')} - Page {slide_number}/{total_slides}"
            footer3.text_frame.paragraphs[0].font.size = Pt(10)
            footer3.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        
        # END OF LOOP - All programmes processed
        
        # === SAVE PPT ===
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        filename = f"Bilan_{partner_name}_{period_label}.pptx".replace(' ', '_').replace('/', '_')
        
        logging.info(f"PPT Generated for {len(programmes)} programmes: {total_slides} slides total")
        
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )
        
    except Exception as e:
        logging.error(f"Error generating PPT: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

# Include the router in the main app
async def export_bilan_partenaire_ppt(
    partenaire_id: str = Query(...),
    date_debut: str = Query(...),
    date_fin: str = Query(...)
):
    """Generate PowerPoint report - STRICT MODE with full validation"""
    try:
        # Logs collection
        logs = {
            "placeholdersReplaced": 0,
            "placeholdersRestants": [],
            "periodUsed": "",
            "periodFallbackApplied": False,
            "rowsSites": 0,
            "rowsLignes": 0,
            "incidentsSynthese": ""
        }
        
        # === GET DATA ===
        partenaire = await db.partenaires.find_one({"id": partenaire_id})
        if not partenaire:
            raise HTTPException(status_code=404, detail="Partenaire not found")
        
        partner_name = partenaire.get('nom', '')
        if not partner_name:
            raise HTTPException(status_code=400, detail="ASSERTION FAIL: partner.name is empty")
        
        # Get programmes
        programme_ids = partenaire.get('programmes_ids', [])
        programmes = await db.programmes.find({"id": {"$in": programme_ids}}).to_list(length=None)
        programmes = sorted(programmes, key=lambda p: p['nom'])
        
        if not programmes:
            raise HTTPException(status_code=404, detail="No programmes found")
        
        programme = programmes[0]
        program_name = programme.get('nom', '')
        if not program_name:
            raise HTTPException(status_code=400, detail="ASSERTION FAIL: program.name is empty")
        
        # === PARSE DATES ===
        try:
            date_debut_obj = datetime.fromisoformat(date_debut).replace(tzinfo=timezone.utc)
            date_fin_obj = datetime.fromisoformat(date_fin).replace(tzinfo=timezone.utc)
            
            # Add one day to date_fin to include the entire day
            date_fin_obj = date_fin_obj.replace(hour=23, minute=59, second=59)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Invalid date format: {str(e)}")
        
        if date_debut_obj > date_fin_obj:
            raise HTTPException(status_code=400, detail="date_debut must be before date_fin")
        
        # Generate period label
        if date_debut_obj.year == date_fin_obj.year and date_debut_obj.month == date_fin_obj.month:
            # Same month
            period_label = f"{format_french_month(date_debut_obj)} {date_debut_obj.year}"
        elif date_debut_obj.year == date_fin_obj.year:
            # Same year, different months
            period_label = f"{format_french_month(date_debut_obj)} - {format_french_month(date_fin_obj)} {date_debut_obj.year}"
        else:
            # Different years
            period_label = f"{format_french_month(date_debut_obj)} {date_debut_obj.year} - {format_french_month(date_fin_obj)} {date_fin_obj.year}"
        
        # Check if period is in future or has no data
        tests_site = await db.tests_site.find({
            "programme_id": programme['id'],
            "partenaire_id": partenaire_id,
            "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
        }).to_list(length=None)
        
        tests_ligne = await db.tests_ligne.find({
            "programme_id": programme['id'],
            "partenaire_id": partenaire_id,
            "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
        }).to_list(length=None)
        
        # FALLBACK: If no data, find last complete month with data
        if len(tests_site) == 0 and len(tests_ligne) == 0:
            logging.warning("PERIOD_FALLBACK: No data for requested period, searching last complete month")
            
            # Find last test date
            last_test_site = await db.tests_site.find({
                "programme_id": programme['id'],
                "partenaire_id": partenaire_id
            }).sort("date_test", -1).to_list(length=1)
            
            last_test_ligne = await db.tests_ligne.find({
                "programme_id": programme['id'],
                "partenaire_id": partenaire_id
            }).sort("date_test", -1).to_list(length=1)
            
            last_date = None
            if last_test_site:
                last_date = datetime.fromisoformat(last_test_site[0]['date_test'])
            if last_test_ligne and (not last_date or datetime.fromisoformat(last_test_ligne[0]['date_test']) > last_date):
                last_date = datetime.fromisoformat(last_test_ligne[0]['date_test'])
            
            if last_date:
                # Use that month
                date_debut_obj = datetime(last_date.year, last_date.month, 1, tzinfo=timezone.utc)
                if last_date.month == 12:
                    date_fin_obj = datetime(last_date.year + 1, 1, 1, tzinfo=timezone.utc)
                else:
                    date_fin_obj = datetime(last_date.year, last_date.month + 1, 1, tzinfo=timezone.utc)
                period_label = f"{format_french_month(date_debut_obj)} {last_date.year}"
                logs["periodFallbackApplied"] = True
                
                # Re-fetch data
                tests_site = await db.tests_site.find({
                    "programme_id": programme['id'],
                    "partenaire_id": partenaire_id,
                    "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
                }).to_list(length=None)
                
                tests_ligne = await db.tests_ligne.find({
                    "programme_id": programme['id'],
                    "partenaire_id": partenaire_id,
                    "date_test": {"$gte": date_debut_obj.isoformat(), "$lt": date_fin_obj.isoformat()}
                }).to_list(length=None)
        
        logs["periodUsed"] = period_label
        
        # Sort by date descending
        tests_site = sorted(tests_site, key=lambda t: t.get('date_test', ''), reverse=True)
        tests_ligne = sorted(tests_ligne, key=lambda t: t.get('date_test', ''), reverse=True)
        
        # Get alertes
        alertes = await db.alertes.find({
            "programme_id": programme['id'],
            "partenaire_id": partenaire_id
        }).to_list(length=None)
        
        # === CALCULATE STATISTICS ===
        total_tests_site = len(tests_site)
        tests_site_reussis = len([t for t in tests_site if t.get('application_remise', False)])
        pct_site = round((tests_site_reussis / total_tests_site * 100), 1) if total_tests_site > 0 else 0
        
        total_tests_ligne = len(tests_ligne)
        tests_ligne_reussis = len([t for t in tests_ligne if t.get('application_offre', False)])
        pct_ligne = round((tests_ligne_reussis / total_tests_ligne * 100), 1) if total_tests_ligne > 0 else 0
        
        # Average waiting time
        delais = []
        for t in tests_ligne:
            if t.get('delai_attente'):
                try:
                    parts = t['delai_attente'].split(':')
                    if len(parts) == 2:
                        delais.append(int(parts[0]) * 60 + int(parts[1]))
                except:
                    pass
        avg_delai = sum(delais) / len(delais) if delais else 0
        avg_delai_str = f"{int(avg_delai // 60):02d}:{int(avg_delai % 60):02d}"
        
        # Average accueil
        accueils = [t.get('evaluation_accueil', '') for t in tests_ligne if t.get('evaluation_accueil')]
        commentaire_accueil = "—"
        if accueils:
            # Count most common
            from collections import Counter
            most_common = Counter(accueils).most_common(1)[0][0]
            commentaire_accueil = most_common
        
        # Partner settings
        naming_remise = partenaire.get('naming_attendu', 'Non communiqué')
        
        # Messagerie/Decroche based on tests ligne
        md_count = len([t for t in tests_ligne if t.get('messagerie_vocale_dediee')])
        dd_count = len([t for t in tests_ligne if t.get('decroche_dedie')])
        info_md = "Oui" if md_count > len(tests_ligne) / 2 else "Non" if len(tests_ligne) > 0 else "—"
        info_dd = "Oui" if dd_count > len(tests_ligne) / 2 else "Non" if len(tests_ligne) > 0 else "—"
        
        # === INCIDENTS SYNTHESE ===
        if alertes:
            nb_incidents = len(alertes)
            types = [inc.get('description', 'N/A')[:30] for inc in alertes]
            from collections import Counter
            top_types = Counter(types).most_common(2)
            synthese_types = ", ".join([t[0] for t in top_types])
            
            logs["incidentsSynthese"] = f"{nb_incidents} alerte(s) détecté(s). Principaux types: {synthese_types}."
        else:
            logs["incidentsSynthese"] = "Aucun alerte sur la période."
        
        # === LOAD TEMPLATE ===
        template_path = TEMPLATE_DIR / "Bilan_Blindtest_template.pptx"
        if not template_path.exists():
            raise HTTPException(status_code=500, detail="Template not found")
        
        work_path = TEMPLATE_DIR / f"temp_{partenaire_id}.pptx"
        shutil.copy(template_path, work_path)
        prs = Presentation(str(work_path))
        
        # === MAPPING DICTIONARY ===
        bilan_date = datetime.now(timezone.utc).strftime('%d/%m/%Y')
        mapping = {
            "PartnerName": partner_name,
            "ProgramName": program_name,
            "Mois du trigger + année du trigger": period_label,
            "moyenne des tests sites réussis": f"{pct_site}%",
            "moyenne des tests lignes réussis": f"{pct_ligne}%",
            "temps d'attente/nombre de test effectués": avg_delai_str,
            "temps d'attente": avg_delai_str,
            "nom de la remise": naming_remise,
            "info case MD": info_md,
            "info case DD": info_dd,
            "commentaire sur l'accueil": commentaire_accueil,
            "Agent qui résume les faits inscrits lors des alertes dans la case dédiée aux commentaires": logs["incidentsSynthese"]
        }
        
        # === REPLACE PLACEHOLDERS WITH TOLERANT REGEX ===
        import re
        placeholder_pattern = re.compile(r'\{\{\s*([A-Za-z0-9_é èàù\'\-]+?)\s*\}\}|\{\s*([A-Za-z0-9_é èàù\'\-]+?)\s*\}', re.IGNORECASE)
        
        def normalize_key(key):
            """Normalize key by removing extra spaces and lowercasing"""
            return key.strip().lower()
        
        # Normalize mapping keys
        normalized_mapping = {normalize_key(k): v for k, v in mapping.items()}
        
        def replace_placeholder(match):
            key = match.group(1) or match.group(2)
            normalized_key = normalize_key(key)
            
            if normalized_key in normalized_mapping:
                logs["placeholdersReplaced"] += 1
                return str(normalized_mapping[normalized_key])
            else:
                logs["placeholdersRestants"].append(f"{{{key}}}")
                return match.group(0)  # Keep original if not mapped
        
        # Replace in all slides
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            original = run.text
                            run.text = placeholder_pattern.sub(replace_placeholder, original)
                
                # Handle tables differently - reconstruct them
                if shape.has_table:
                    table = shape.table
                    
                    # Determine table type from slide content
                    is_sites_table = False
                    is_lignes_table = False
                    
                    for s in slide.shapes:
                        if s.has_text_frame and 'Sites' in s.text and 'Ligne' not in s.text:
                            is_sites_table = True
                        if s.has_text_frame and 'Ligne' in s.text:
                            is_lignes_table = True
                    
                    if is_sites_table:
                        # Fill Sites table
                        if tests_site:
                            for test in tests_site:
                                try:
                                    test_date = datetime.fromisoformat(test['date_test'])
                                    row_data = [
                                        test_date.strftime('%d/%m/%Y'),
                                        test.get('url', 'N/A'),
                                        test.get('url', 'N/A'),
                                        'OK' if test.get('application_remise') else 'KO',
                                        'N/A',
                                        str(test.get('pct_remise_calcule', 0)),
                                        test.get('naming_constate', '')
                                    ]
                                    fill_table_with_data(table, [row_data], header_rows=1)
                                    logs["rowsSites"] += 1
                                except Exception as e:
                                    logging.error(f"Error filling sites table: {str(e)}")
                        else:
                            fill_table_with_data(table, [["Aucun test disponible pour cette période", "", "", "", "", "", ""]], header_rows=1)
                    
                    elif is_lignes_table:
                        # Fill Lignes table
                        if tests_ligne:
                            for test in tests_ligne:
                                try:
                                    test_date = datetime.fromisoformat(test['date_test'])
                                    row_data = [
                                        test_date.strftime('%d/%m/%Y'),
                                        test.get('numero_telephone', 'N/A'),
                                        'OK' if test.get('application_offre') else 'KO',
                                        test.get('delai_attente', 'N/A'),
                                        'Oui' if test.get('messagerie_vocale_dediee') else 'Non',
                                        'Oui' if test.get('decroche_dedie') else 'Non',
                                        test.get('evaluation_accueil', 'N/A'),
                                        ''
                                    ]
                                    fill_table_with_data(table, [row_data], header_rows=1)
                                    logs["rowsLignes"] += 1
                                except Exception as e:
                                    logging.error(f"Error filling lignes table: {str(e)}")
                        else:
                            fill_table_with_data(table, [["Aucun test disponible pour cette période", "", "", "", "", "", "", ""]], header_rows=1)
        
        # === FINAL ASSERTION: Check for remaining placeholders ===
        remaining_placeholders = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text
                    matches = re.findall(r'\{[^}]+\}', text)
                    if matches:
                        for match in matches:
                            remaining_placeholders.append(match)
        
        if remaining_placeholders and len(logs['placeholdersRestants']) > 0:
            raise HTTPException(
                status_code=500,
                detail=f"ASSERTION FAIL: Unresolved placeholders found: {list(set(remaining_placeholders))}"
            )
        
        # === CHECK TABLE REBUILD ===
        if (len(tests_site) > 0 or len(tests_ligne) > 0) and (logs["rowsSites"] == 0 and logs["rowsLignes"] == 0):
            raise HTTPException(status_code=500, detail="TABLES_REBUILD_FAILED: Data exists but no rows inserted")
        
        # === SAVE PPT ===
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        # Clean up
        try:
            work_path.unlink()
        except:
            pass
        
        filename = f"Bilan_Blindtest_{partner_name}_{program_name}_{period_label}.pptx".replace(' ', '_').replace('/', '_')
        
        # Log summary
        logging.info(f"PPT Generated: {logs}")
        
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )
        
    except Exception as e:
        logging.error(f"Error generating PPT: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

# Include the router in the main app
app.include_router(api_router)

# Mount static files for uploads
app.mount("/uploads", StaticFiles(directory=str(UPLOAD_DIR)), name="uploads")

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()